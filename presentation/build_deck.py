"""
BearWatch BSI - research presentation deck v9 (cream + orange edition).

Visual redesign from v8:
  - Cream / off-white background instead of dark ink
  - Illinois orange as primary accent throughout
  - Aggressive text reduction: each slide = one headline + one visual + minimal copy
  - Trust the verbal narrative; the deck should NOT be a teleprompter

Same 19-slide structure as v8 (Title -> Q&A) but rebuilt for clarity and visual breath.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path("C:/Users/siddh/Desktop/spring 2026/580/BNPL_v9_FINAL/03_presentation/BearWatch_Research_Deck_v9o.pptx")
LOGO = Path("C:/Users/siddh/Desktop/Fall 2025/Illinois-Logo-No-background.png")

SCREENSHOT_DIR = Path("C:/Users/siddh/Desktop/spring 2026/580/BNPL_v9_FINAL/03_presentation/screenshots")
CVNA_CHART_PNG    = SCREENSHOT_DIR / "cvna_chart.png"     # image 1 — SHORT trade card + backtest replay
CVNA_MATH_PNG     = SCREENSHOT_DIR / "cvna_math.png"      # image 2 — "The trade math" table
CVNA_PLAYBOOK_PNG = SCREENSHOT_DIR / "cvna_playbook.png"  # image 3 — Updated Trade Playbook · both sides
LIVE_POD_PNG      = SCREENSHOT_DIR / "live_pod.png"

# Backwards-compat alias (slide 17 expects this name)
CASE_STUDY_PNG = CVNA_CHART_PNG

# === CREAM + ORANGE PALETTE ===================================================
# Background family — warm cream
BG          = RGBColor(0xFA, 0xF6, 0xEC)   # primary canvas
BG_PANEL    = RGBColor(0xF2, 0xEB, 0xD9)   # panel fill (slightly darker cream)
BG_DEEP     = RGBColor(0xE9, 0xDF, 0xC4)   # deeper accent panel
HAIRLINE    = RGBColor(0xD8, 0xCB, 0xAA)   # warm divider
HAIRLINE_2  = RGBColor(0xC0, 0xB1, 0x8E)   # stronger divider

# Type
TEXT        = RGBColor(0x1F, 0x1B, 0x16)   # near-black, warm
TEXT_DIM    = RGBColor(0x5A, 0x52, 0x47)   # warm gray secondary
TEXT_FAINT  = RGBColor(0x8C, 0x84, 0x76)   # subtle tertiary

# Illinois orange — primary accent
ORANGE      = RGBColor(0xE8, 0x4A, 0x27)   # core
ORANGE_DEEP = RGBColor(0xB8, 0x35, 0x18)   # darker for borders / serious moments
ORANGE_SOFT = RGBColor(0xF4, 0xC0, 0xA8)   # soft fill for backgrounds

# Restrained secondary accents (matched to warm canvas)
GOLD        = RGBColor(0xB8, 0x84, 0x1C)   # honest-disclosure / warnings
GREEN       = RGBColor(0x2D, 0x7A, 0x4F)   # validated / success
RED_TERRA   = RGBColor(0xB8, 0x56, 0x38)   # nulls / risks (muted terracotta)
SLATE       = RGBColor(0x4A, 0x6E, 0x8F)   # neutral data accent
PLUM        = RGBColor(0x6B, 0x4A, 0x7F)   # operational / pod accent

FONT_DISP = "Cambria"           # calligraphic serif — distinctive, refined, academic
FONT_HEAD = "Segoe UI"          # clean sans-serif for body
FONT_MONO = "Consolas"          # data tables / acronyms

# === Section colors (slide-number → (label, accent)) ===========================
# Used by kicker() to display a section pill above each slide's heading.
# Sections: Foundation (slate) → Methodology (orange) → Findings (gold)
#           → Validation (green) → Direction (plum) → Pod (deep orange)
SECTION_OF = {
    1:  ("FOUNDATION",          "#4A6E8F"),
    2:  ("FOUNDATION",          "#4A6E8F"),
    3:  ("FOUNDATION",          "#4A6E8F"),
    4:  ("FOUNDATION",          "#4A6E8F"),
    5:  ("FOUNDATION",          "#4A6E8F"),
    6:  ("FOUNDATION",          "#4A6E8F"),
    7:  ("METHODOLOGY",         "#E84A27"),
    8:  ("METHODOLOGY",         "#E84A27"),
    9:  ("METHODOLOGY",         "#E84A27"),
    10: ("METHODOLOGY",         "#E84A27"),
    11: ("EMPIRICAL FINDINGS",  "#B8841C"),
    12: ("EMPIRICAL FINDINGS",  "#B8841C"),
    13: ("EMPIRICAL FINDINGS",  "#B8841C"),
    14: ("EMPIRICAL FINDINGS",  "#B8841C"),
    15: ("EMPIRICAL FINDINGS",  "#B8841C"),
    16: ("VALIDATION & SCOPE",  "#2D7A4F"),
    17: ("VALIDATION & SCOPE",  "#2D7A4F"),
    18: ("VALIDATION & SCOPE",  "#2D7A4F"),
    19: ("VALIDATION & SCOPE",  "#2D7A4F"),
    20: ("VALIDATION & SCOPE",  "#2D7A4F"),
    21: ("VALIDATION & SCOPE",  "#2D7A4F"),
    22: ("RESEARCH DIRECTION",  "#6B4A7F"),
    23: ("RESEARCH DIRECTION",  "#6B4A7F"),
    24: ("RESEARCH DIRECTION",  "#6B4A7F"),
    25: ("OPERATIONAL POD",     "#B83518"),
    26: ("OPERATIONAL POD",     "#B83518"),
    27: ("OPERATIONAL POD",     "#B83518"),
}

# === Setup ====================================================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
M = Inches(0.85)

TOTAL = 28


# ===============================================================================
# Helpers — light-theme rebuild
# ===============================================================================
def _rgb_hex(c):
    """RGBColor → uppercase 6-digit hex string for OOXML."""
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def bg(slide, color=BG):
    """Gradient background (warm cream top → slightly deeper cream bottom).
    Falls back to solid fill if XML injection isn't supported."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.line.fill.background()
    r.shadow.inherit = False
    # Inject gradient fill via OOXML to get a true two-stop gradient
    try:
        from lxml import etree
        # spPr is the shape's <p:spPr> element
        spPr = r.fill._xPr
        # Strip any existing solid/grad/no fill children
        for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:pattFill"):
            for el in spPr.findall(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag.split(':')[1]}"):
                spPr.remove(el)
        c1 = _rgb_hex(BG)         # top — light warm cream
        c2 = _rgb_hex(BG_DEEP)    # bottom — slightly deeper cream
        grad_xml = (
            '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">'
              '<a:gsLst>'
                f'<a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>'
                f'<a:gs pos="100000"><a:srgbClr val="{c2}"/></a:gs>'
              '</a:gsLst>'
              '<a:lin ang="5400000" scaled="0"/>'  # 90° = top-to-bottom
            '</a:gradFill>'
        )
        # Insert before the line element if present, else just append
        from copy import deepcopy
        elem = etree.fromstring(grad_xml)
        # Find <a:ln> child to insert before (so fill comes before line in OOXML order)
        ln_tag = "{http://schemas.openxmlformats.org/drawingml/2006/main}ln"
        ln = spPr.find(ln_tag)
        if ln is not None:
            ln.addprevious(elem)
        else:
            spPr.append(elem)
    except Exception:
        # Solid fallback
        r.fill.solid()
        r.fill.fore_color.rgb = color
    return r


def text(slide, s, x, y, w, h, *, size=14, color=TEXT, bold=False,
         align=PP_ALIGN.LEFT, font=FONT_HEAD, anchor=MSO_ANCHOR.TOP, italic=False,
         line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing: p.line_spacing = line_spacing
    r = p.add_run()
    r.text = s
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


def block(slide, paragraphs, x, y, w, h, *, line_spacing=1.4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, pd in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pd.get("align", PP_ALIGN.LEFT)
        p.line_spacing = pd.get("line_spacing", line_spacing)
        if pd.get("space_before"): p.space_before = Pt(pd["space_before"])
        r = p.add_run()
        r.text = pd["text"]
        r.font.name = pd.get("font", FONT_HEAD)
        r.font.size = Pt(pd.get("size", 14))
        r.font.color.rgb = pd.get("color", TEXT)
        r.font.bold = pd.get("bold", False)
        r.font.italic = pd.get("italic", False)
    return tb


def panel(slide, x, y, w, h, *, fill=BG_PANEL, border=HAIRLINE, border_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.03
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(border_w)
    s.shadow.inherit = False
    return s


def hairline(slide, x1, y1, x2, y2, color=HAIRLINE, w=0.75):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(w)
    return line


def orange_bar(slide, x, y, w, h):
    """Solid orange accent bar — used as a vertical rule."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.line.fill.background()
    r.fill.solid()
    r.fill.fore_color.rgb = ORANGE
    r.shadow.inherit = False
    return r


def kicker(slide, num, label):
    """Top-left section marker — section pill + slide num + label."""
    # Resolve which broader section this slide belongs to
    section_label, section_hex = SECTION_OF.get(num, ("", "#E84A27"))
    section_rgb = RGBColor(int(section_hex[1:3], 16),
                           int(section_hex[3:5], 16),
                           int(section_hex[5:7], 16))

    # Section name (small mono caps, top-left)
    if section_label:
        text(slide, section_label, M, Inches(0.38), Inches(4.0), Inches(0.25),
             size=9, color=section_rgb, bold=True, font=FONT_MONO)

    # Slide num + slide label (right-aligned for slide num, left-aligned for label)
    text(slide, f"{num:02d}", M, Inches(0.65), Inches(0.55), Inches(0.4),
         size=15, color=section_rgb, bold=True, font=FONT_MONO)
    text(slide, label.upper(), M + Inches(0.65), Inches(0.68), SW - 2 * M - Inches(0.65), Inches(0.36),
         size=12, color=TEXT, bold=True, font=FONT_MONO)

    # Section-colored hairline (slightly thicker, tighter to type)
    hairline(slide, M, Inches(1.05), SW - M, Inches(1.05), section_rgb, w=1.4)


def footer(slide, page_num):
    text(slide, "BearWatch BSI  ·  FIN 580  ·  Spring 2026",
         M, SH - Inches(0.45), Inches(7), Inches(0.3),
         size=9, color=TEXT_FAINT, font=FONT_MONO)
    text(slide, f"{page_num:02d} / {TOTAL:02d}",
         SW - Inches(2.0), SH - Inches(0.45), Inches(0.8), Inches(0.3),
         size=9, color=TEXT_FAINT, font=FONT_MONO, align=PP_ALIGN.RIGHT)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), SW - Inches(1.05), SH - Inches(0.55),
                                 height=Inches(0.4))


def slide_title(slide, title, sub=None):
    """Big calligraphic headline + clean subtitle in dark warm — no orange clash."""
    text(slide, title, M, Inches(1.25), SW - 2 * M, Inches(0.95),
         size=38, color=TEXT, bold=True, font=FONT_DISP, line_spacing=1.1)
    if sub:
        text(slide, sub, M, Inches(2.20), SW - 2 * M, Inches(0.55),
             size=15, color=TEXT_DIM, bold=True, italic=False, font=FONT_HEAD)


# ===============================================================================
# SLIDE 1 - Title
# ===============================================================================
def s_title():
    s = prs.slides.add_slide(BLANK); bg(s)

    # Logo top-centre
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), SW / 2 - Inches(0.65), Inches(0.7),
                             height=Inches(1.05))

    # Orange divider
    orange_bar(s, M, Inches(2.15), SW - 2 * M, Inches(0.04))

    # Big title — BOLD
    text(s, "Alternative-Data Leading Indicators",
         M, Inches(2.55), SW - 2 * M, Inches(1.0),
         size=48, color=TEXT, bold=True, font=FONT_DISP, align=PP_ALIGN.CENTER)
    text(s, "of Consumer-Credit Distress",
         M, Inches(3.4), SW - 2 * M, Inches(1.0),
         size=48, color=ORANGE_DEEP, bold=True, font=FONT_DISP, align=PP_ALIGN.CENTER)

    # Working-paper kicker
    text(s, "WORKING PAPER  ·  EMPIRICAL FINANCE",
         M, Inches(1.95), SW - 2 * M, Inches(0.3),
         size=12, color=ORANGE, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # Subtitle one-liner
    text(s, "An empirical investigation of the Behavioural Stress Index across 27 consumer-credit issuers, 2019–2026.",
         M, Inches(4.5), SW - 2 * M, Inches(0.6),
         size=16, color=TEXT, bold=True, italic=False, font=FONT_HEAD, align=PP_ALIGN.CENTER)

    # Two-row tagline strip — umbrella term + the universe
    panel(s, M + Inches(0.5), Inches(5.2), SW - 2 * M - Inches(1.0), Inches(1.0),
          fill=ORANGE_SOFT, border=ORANGE, border_w=1.2)
    text(s, "UNSECURED NON-BANK CONSUMER CREDIT  ·  installment lending / Buy Now, Pay Later (BNPL)",
         M + Inches(0.5), Inches(5.32), SW - 2 * M - Inches(1.0), Inches(0.4),
         size=13, color=ORANGE_DEEP, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    text(s, "BNPL  ·  Subprime Auto  ·  Marketplace Lending  ·  Consumer Finance  ·  Card Monolines",
         M + Inches(0.5), Inches(5.75), SW - 2 * M - Inches(1.0), Inches(0.4),
         size=12, color=ORANGE_DEEP, bold=True, italic=False, align=PP_ALIGN.CENTER)

    # Bottom rule + meta
    hairline(s, M, SH - Inches(1.3), SW - M, SH - Inches(1.3), HAIRLINE_2, w=1.2)
    block(s, [
        {"text": "Siddharth Verma",
         "size": 16, "color": TEXT, "bold": True, "align": PP_ALIGN.CENTER},
        {"text": "FIN 580  ·  University of Illinois Urbana-Champaign  ·  Spring 2026",
         "size": 12, "color": TEXT_DIM, "bold": True, "font": FONT_MONO, "align": PP_ALIGN.CENTER},
    ], M, SH - Inches(1.05), SW - 2 * M, Inches(0.8))


# ===============================================================================
# SLIDE 2 - Research motivation (lit gap)
# ===============================================================================
def s_motivation():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 1, "Research motivation")
    slide_title(s, "An untested application of an established literature.",
                "CFPB-as-distress-predictor exists for banks — not for unsecured non-bank consumer credit (installment lending / BNPL).")

    # === BNPL stats strip — four facts ===
    stat_y = Inches(2.95)
    stat_h = Inches(1.05)
    stat_w = (SW - 2 * M - Inches(0.6)) / 4
    stats = [
        ("2.4×",       "more likely subprime",    "TransUnion 2024"),
        ("10.5%",      "BNPL late-payment rate",  "CFPB 2021"),
        ("34→39→42%",  "industry late-pay traj.", "2023 → 2024 → 2025e"),
        ("Nov 2023",   "Fed flagged segment",     "Financial Stability Report"),
    ]
    for i, (big, mid, src) in enumerate(stats):
        x = M + i * (stat_w + Inches(0.2))
        panel(s, x, stat_y, stat_w, stat_h,
              fill=ORANGE_SOFT, border=ORANGE, border_w=1.2)
        text(s, big, x, stat_y + Inches(0.1), stat_w, Inches(0.5),
             size=24, color=ORANGE_DEEP, bold=True, font=FONT_DISP, align=PP_ALIGN.CENTER)
        text(s, mid, x + Inches(0.1), stat_y + Inches(0.55), stat_w - Inches(0.2), Inches(0.3),
             size=11, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        text(s, src, x + Inches(0.1), stat_y + Inches(0.78), stat_w - Inches(0.2), Inches(0.3),
             size=9, color=TEXT_DIM, italic=False, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # === Two columns below: established work | open question ===
    body_y = Inches(4.25)
    body_h = Inches(2.4)
    col_w = (SW - 2 * M - Inches(0.4)) / 2

    # LEFT — established
    panel(s, M, body_y, col_w, body_h)
    text(s, "ESTABLISHED  ·  CFPB → bank distress",
         M + Inches(0.3), body_y + Inches(0.12), col_w, Inches(0.3),
         size=11, color=GREEN, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "Begley & Purnanandam (2021, RFS)",
         "size": 13, "color": TEXT, "bold": True, "space_before": 4},
        {"text": "Complaint quality predicts regulatory action.",
         "size": 12, "color": TEXT_DIM, "bold": True, "italic": False, "space_before": 2},

        {"text": "Dou & Wang (2023)",
         "size": 13, "color": TEXT, "bold": True, "space_before": 8},
        {"text": "Complaint volume predicts bank loan losses 2-4 q ahead.",
         "size": 12, "color": TEXT_DIM, "bold": True, "italic": False, "space_before": 2},

        {"text": "Hayes, Jiang & Pan (2021)",
         "size": 13, "color": TEXT, "bold": True, "space_before": 8},
        {"text": "Complaint metadata predicts supervisory outcomes.",
         "size": 12, "color": TEXT_DIM, "bold": True, "italic": False, "space_before": 2},
    ], M + Inches(0.3), body_y + Inches(0.5), col_w - Inches(0.6), body_h - Inches(0.6))

    # RIGHT — open question
    panel(s, M + col_w + Inches(0.4), body_y, col_w, body_h,
          border=ORANGE, border_w=1.5)
    text(s, "OPEN QUESTION  ·  this paper",
         M + col_w + Inches(0.7), body_y + Inches(0.12), col_w, Inches(0.3),
         size=11, color=ORANGE, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "Does the result generalise to non-bank consumer credit?",
         "size": 15, "color": TEXT, "bold": True, "space_before": 4, "line_spacing": 1.35},
        {"text": "Methodological gap — no prior multi-pillar fusion.",
         "size": 12, "color": TEXT_DIM, "bold": True, "italic": False, "space_before": 12, "line_spacing": 1.5},
        {"text": "Published CFPB work uses one alt-data source at a time. We fuse CFPB with seven additional pillars under pre-registered weights.",
         "size": 11, "color": TEXT_DIM, "bold": True, "italic": False, "space_before": 6, "line_spacing": 1.5},
    ], M + col_w + Inches(0.7), body_y + Inches(0.5), col_w - Inches(0.6), body_h - Inches(0.6))

    footer(s, 2)


# ===============================================================================
# SLIDE 3 - What is Behavioural Stress Index (BSI)? (NEW · introduce the instrument visually)
# ===============================================================================
def s_what_is_bsi():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 2, "What is BSI?")
    slide_title(s, "Behavioural Stress Index  =  one number per firm per day.",
                "An 8-pillar composite z-score that detects spikes in customer-distress volume.")

    # Three-column flow:  INPUT  →  ENGINE  →  OUTPUT
    col_y = Inches(2.95)
    col_h = Inches(3.4)
    in_w   = Inches(4.2)
    eng_w  = Inches(3.6)
    out_w  = SW - 2 * M - in_w - eng_w - Inches(0.5)
    in_x   = M
    eng_x  = in_x + in_w + Inches(0.25)
    out_x  = eng_x + eng_w + Inches(0.25)

    # === LEFT — INPUTS ===
    panel(s, in_x, col_y, in_w, col_h)
    text(s, "INPUT  ·  8 PUBLIC ALT-DATA PILLARS",
         in_x + Inches(0.25), col_y + Inches(0.15), in_w, Inches(0.3),
         size=11, color=ORANGE, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "·  Consumer Financial Protection Bureau (CFPB) complaints   (weight 0.45 ★)",
         "size": 12, "color": TEXT, "bold": True, "space_before": 8, "font": FONT_MONO},
        {"text": "·  App-store reviews   (0.15)",
         "size": 12, "color": TEXT_DIM, "bold": True, "space_before": 4, "font": FONT_MONO},
        {"text": "·  Reddit text   (0.10)",
         "size": 12, "color": TEXT_DIM, "bold": True, "space_before": 4, "font": FONT_MONO},
        {"text": "·  Bluesky text   (0.05)",
         "size": 12, "color": TEXT_DIM, "bold": True, "space_before": 4, "font": FONT_MONO},
        {"text": "·  Google Trends   (0.10)",
         "size": 12, "color": TEXT_DIM, "bold": True, "space_before": 4, "font": FONT_MONO},
        {"text": "·  ABS pool metrics   (0.10)",
         "size": 12, "color": TEXT_DIM, "bold": True, "space_before": 4, "font": FONT_MONO},
        {"text": "·  8-K filings   (0.05)",
         "size": 12, "color": TEXT_DIM, "bold": True, "space_before": 4, "font": FONT_MONO},
        {"text": "·  Merrill Option Volatility Estimate index (MOVE) / FRED macro   (gate G3)",
         "size": 12, "color": TEXT_DIM, "bold": True, "space_before": 4, "font": FONT_MONO},
        {"text": "7 of 8 are volume-derived series.",
         "size": 11, "color": ORANGE_DEEP, "bold": True, "italic": False, "space_before": 14},
    ], in_x + Inches(0.25), col_y + Inches(0.5), in_w - Inches(0.5), col_h - Inches(0.6))

    # Arrow 1
    text(s, "→", in_x + in_w + Inches(0.02), col_y + Inches(1.5), Inches(0.25), Inches(0.5),
         size=24, color=ORANGE, bold=True, font=FONT_DISP, align=PP_ALIGN.CENTER)

    # === CENTER — ENGINE ===
    panel(s, eng_x, col_y, eng_w, col_h,
          fill=ORANGE_SOFT, border=ORANGE, border_w=1.5)
    text(s, "ENGINE  ·  EWMA z-score fusion",
         eng_x + Inches(0.2), col_y + Inches(0.15), eng_w, Inches(0.3),
         size=11, color=ORANGE_DEEP, bold=True, font=FONT_MONO)

    text(s, "BSIₜ  =  Σᵢ  γᵢ · wᵢ · z^EWMA(Xᵢ,ₜ)",
         eng_x + Inches(0.1), col_y + Inches(0.95), eng_w - Inches(0.2), Inches(0.6),
         size=18, color=TEXT, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)

    block(s, [
        {"text": "z-score each pillar against",
         "size": 12, "color": TEXT, "bold": True, "space_before": 0, "align": PP_ALIGN.CENTER, "line_spacing": 1.3},
        {"text": "its own running history",
         "size": 12, "color": TEXT, "bold": True, "space_before": 0, "align": PP_ALIGN.CENTER, "line_spacing": 1.3},
        {"text": "(EWMA half-life 250d)",
         "size": 11, "color": ORANGE_DEEP, "bold": True, "italic": False, "space_before": 4, "font": FONT_MONO, "align": PP_ALIGN.CENTER},
        {"text": "Coverage gates zero out missing.",
         "size": 11, "color": TEXT_DIM, "bold": True, "italic": False, "space_before": 14, "align": PP_ALIGN.CENTER, "line_spacing": 1.4},
        {"text": "Pre-registered weights · hash-locked.",
         "size": 11, "color": TEXT_DIM, "bold": True, "italic": False, "space_before": 4, "align": PP_ALIGN.CENTER, "line_spacing": 1.4},
    ], eng_x + Inches(0.15), col_y + Inches(1.7), eng_w - Inches(0.3), Inches(1.6))

    # Arrow 2
    text(s, "→", eng_x + eng_w + Inches(0.02), col_y + Inches(1.5), Inches(0.25), Inches(0.5),
         size=24, color=ORANGE, bold=True, font=FONT_DISP, align=PP_ALIGN.CENTER)

    # === RIGHT — OUTPUT ===
    panel(s, out_x, col_y, out_w, col_h)
    text(s, "OUTPUT  ·  one z-score per firm per day",
         out_x + Inches(0.2), col_y + Inches(0.15), out_w, Inches(0.3),
         size=11, color=GREEN, bold=True, font=FONT_MONO)

    # Mini "thermometer" thresholds
    bands_y = col_y + Inches(0.6)
    bands = [
        ("z ≥ 2.5",  "GUARDIAN fires",                 GREEN),
        ("z ≥ 2.0",  "SCOUT fires  ★",                 ORANGE),
        ("z ≥ 1.5",  "WORRIED state",                  GOLD),
        ("|z| < 0.5", "SLEEPING — benign",              TEXT_DIM),
        ("z negative", "below baseline — quiet",         TEXT_FAINT),
    ]
    bh = Inches(0.4)
    for i, (z, lbl, col) in enumerate(bands):
        y = bands_y + i * bh
        # color rule
        rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  out_x + Inches(0.25), y + Inches(0.08),
                                  Inches(0.05), Inches(0.28))
        rule.line.fill.background(); rule.fill.solid()
        rule.fill.fore_color.rgb = col; rule.shadow.inherit = False
        text(s, z, out_x + Inches(0.4), y, Inches(1.3), bh,
             size=11, color=col, bold=True, font=FONT_MONO)
        text(s, lbl, out_x + Inches(1.55), y, out_w - Inches(1.7), bh,
             size=11, color=TEXT_DIM, bold=True)

    # Bottom punch line
    text(s, "Cross-firm comparable.  Real-time.  Pre-registered.  Replayable byte-for-byte.",
         out_x + Inches(0.25), col_y + col_h - Inches(0.45), out_w - Inches(0.5), Inches(0.35),
         size=10, color=ORANGE_DEEP, bold=True, italic=False, font=FONT_MONO, line_spacing=1.4)

    # Bottom strip
    panel(s, M, Inches(6.45), SW - 2 * M, Inches(0.55),
          fill=ORANGE_SOFT, border=ORANGE, border_w=1.0)
    text(s, "BSI is a leading indicator on changes in default probability — not a default model itself.",
         M + Inches(0.4), Inches(6.55), SW - 2 * M - Inches(0.8), Inches(0.4),
         size=13, color=ORANGE_DEEP, bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER)
    footer(s, 3)


# ===============================================================================
# SLIDE 4 - 2D NEURAL DIAGRAM (architecture)
# ===============================================================================
def s_neural_diagram():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 3, "BSI architecture · 2D node diagram")
    slide_title(s, "Inputs → engine → gates → trade verdict.",
                "Same shape as a feed-forward network — every connection is deterministic and pre-registered.")

    # Layer geometry
    diagram_top = Inches(2.85)
    diagram_h   = Inches(3.8)

    # 4 columns: inputs (8 nodes) | engine (1 node) | gates (5 nodes) | verdict (1 node)
    col1_x = M + Inches(0.3)            # inputs
    col2_x = M + Inches(3.5)            # engine
    col3_x = M + Inches(6.6)            # gates
    col4_x = M + Inches(9.7)            # verdict

    node_w_in    = Inches(2.1)
    node_h_in    = Inches(0.32)
    engine_w     = Inches(2.7)
    engine_h     = Inches(2.6)
    gate_w       = Inches(2.1)
    gate_h       = Inches(0.45)
    verdict_w    = Inches(2.1)
    verdict_h    = Inches(2.6)

    # Layer headers
    headers = [
        (col1_x,  Inches(2.6), node_w_in, "INPUT  ·  8 PILLARS",      "#4A6E8F"),
        (col2_x,  Inches(2.6), engine_w,  "ENGINE  ·  exponentially-weighted moving average (EWMA) + WEIGHTS","#E84A27"),
        (col3_x,  Inches(2.6), gate_w,    "GATES  ·  G1 → G5",        "#B8841C"),
        (col4_x,  Inches(2.6), verdict_w, "VERDICT",                   "#2D7A4F"),
    ]
    for hx, hy, hw, htxt, hex_col in headers:
        col = RGBColor(int(hex_col[1:3], 16), int(hex_col[3:5], 16), int(hex_col[5:7], 16))
        text(s, htxt, hx, hy, hw, Inches(0.25),
             size=10, color=col, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # Input nodes — 8 pillars
    pillars = [
        ("CFPB",          "0.45 ★", "#E84A27"),
        ("App-store",     "0.15",   SLATE),
        ("Reddit",        "0.10",   SLATE),
        ("Bluesky",       "0.05",   SLATE),
        ("Trends",        "0.10",   SLATE),
        ("ABS pool",      "0.10",   SLATE),
        ("8-K filings",   "0.05",   SLATE),
        ("MOVE / FRED",   "G3",     "#6B4A7F"),
    ]
    pillar_y_positions = []
    n_pillars = len(pillars)
    spacing = (diagram_h - Inches(0.2)) / (n_pillars - 1)
    for i, (name, weight, color_hex) in enumerate(pillars):
        py = diagram_top + i * spacing
        pillar_y_positions.append(py + node_h_in / 2)
        col = (color_hex if isinstance(color_hex, RGBColor)
               else RGBColor(int(color_hex[1:3], 16), int(color_hex[3:5], 16), int(color_hex[5:7], 16)))
        # Node
        node = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  col1_x, py, node_w_in, node_h_in)
        node.adjustments[0] = 0.5
        node.fill.solid(); node.fill.fore_color.rgb = BG_PANEL
        node.line.color.rgb = col; node.line.width = Pt(1.0)
        node.shadow.inherit = False
        # Label
        text(s, name, col1_x + Inches(0.1), py + Inches(0.04),
             node_w_in - Inches(0.85), Inches(0.25),
             size=10, color=TEXT, bold=True)
        text(s, weight, col1_x + node_w_in - Inches(0.7), py + Inches(0.04),
             Inches(0.6), Inches(0.25),
             size=9, color=col, bold=True, font=FONT_MONO, align=PP_ALIGN.RIGHT)

    # Engine node — large rounded rect
    engine_y = diagram_top + (diagram_h - engine_h) / 2
    engine = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                col2_x, engine_y, engine_w, engine_h)
    engine.adjustments[0] = 0.08
    engine.fill.solid(); engine.fill.fore_color.rgb = ORANGE_SOFT
    engine.line.color.rgb = ORANGE; engine.line.width = Pt(2.0)
    engine.shadow.inherit = False
    text(s, "BSI ENGINE",
         col2_x, engine_y + Inches(0.2), engine_w, Inches(0.4),
         size=14, color=ORANGE_DEEP, bold=True, font=FONT_DISP, align=PP_ALIGN.CENTER)
    text(s, "weighted EWMA",
         col2_x, engine_y + Inches(0.65), engine_w, Inches(0.3),
         size=10, color=ORANGE_DEEP, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    text(s, "z-score fusion",
         col2_x, engine_y + Inches(0.92), engine_w, Inches(0.3),
         size=10, color=ORANGE_DEEP, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    text(s, "BSIₜ  =  Σᵢ γᵢ wᵢ z(Xᵢ,ₜ)",
         col2_x, engine_y + Inches(1.45), engine_w, Inches(0.4),
         size=11, color=TEXT, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    text(s, "h.l. = 250 d   ·   Σwᵢ = 1",
         col2_x, engine_y + Inches(1.85), engine_w, Inches(0.3),
         size=9, color=TEXT_DIM, italic=False, font=FONT_MONO, align=PP_ALIGN.CENTER)
    text(s, "→  one z per firm/day",
         col2_x, engine_y + Inches(2.18), engine_w, Inches(0.3),
         size=9, color=ORANGE_DEEP, bold=True, italic=False, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # Gate nodes — 5 stacked
    gates = [
        ("G1  BSI",   "z ≥ 2.0",     ORANGE),
        ("G2  Subprime-Credit Profile (SCP)",   "phase 2",     SLATE),
        ("G3  MOVE",  "z ≤ 1.0",     PLUM),
        ("G4  Consumer-Credit Divergence (CCD)",   "elevated",    GOLD),
        ("G5  Fundamentals Distress Score (FDS)",   "EDGAR XBRL",  GREEN),
    ]
    gate_y_positions = []
    g_spacing = (diagram_h - Inches(0.2)) / (len(gates) - 1)
    for i, (gname, gcond, gcol) in enumerate(gates):
        gy = diagram_top + i * g_spacing
        gate_y_positions.append(gy + gate_h / 2)
        node = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  col3_x, gy, gate_w, gate_h)
        node.adjustments[0] = 0.4
        node.fill.solid(); node.fill.fore_color.rgb = BG_PANEL
        node.line.color.rgb = gcol; node.line.width = Pt(1.5)
        node.shadow.inherit = False
        text(s, gname, col3_x + Inches(0.12), gy + Inches(0.06),
             gate_w / 2, Inches(0.32),
             size=11, color=TEXT, bold=True, font=FONT_MONO)
        text(s, gcond, col3_x + gate_w / 2, gy + Inches(0.07),
             gate_w / 2 - Inches(0.12), Inches(0.32),
             size=10, color=gcol, bold=True, italic=False, font=FONT_MONO,
             align=PP_ALIGN.RIGHT)

    # Verdict node — large rounded rect on right
    verdict_y = diagram_top + (diagram_h - verdict_h) / 2
    verdict = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 col4_x, verdict_y, verdict_w, verdict_h)
    verdict.adjustments[0] = 0.08
    verdict.fill.solid(); verdict.fill.fore_color.rgb = BG_PANEL
    verdict.line.color.rgb = GREEN; verdict.line.width = Pt(2.0)
    verdict.shadow.inherit = False
    text(s, "TRADE",
         col4_x, verdict_y + Inches(0.25), verdict_w, Inches(0.4),
         size=15, color=GREEN, bold=True, font=FONT_DISP, align=PP_ALIGN.CENTER)
    text(s, "VERDICT",
         col4_x, verdict_y + Inches(0.7), verdict_w, Inches(0.4),
         size=15, color=GREEN, bold=True, font=FONT_DISP, align=PP_ALIGN.CENTER)
    text(s, "APPROVED",
         col4_x, verdict_y + Inches(1.3), verdict_w, Inches(0.3),
         size=10, color=GREEN, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    text(s, "SCALED 40%",
         col4_x, verdict_y + Inches(1.62), verdict_w, Inches(0.3),
         size=10, color=GOLD, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    text(s, "BLOCKED",
         col4_x, verdict_y + Inches(1.94), verdict_w, Inches(0.3),
         size=10, color=RED_TERRA, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # === Connections ===
    # Inputs → Engine (8 lines converging)
    engine_in_y = engine_y + engine_h / 2
    for py in pillar_y_positions:
        ln = s.shapes.add_connector(1, col1_x + node_w_in, py,
                                       col2_x, engine_in_y)
        ln.line.color.rgb = HAIRLINE_2; ln.line.width = Pt(0.6)

    # Engine → Gates (5 lines diverging)
    for gy in gate_y_positions:
        ln = s.shapes.add_connector(1, col2_x + engine_w, engine_in_y,
                                       col3_x, gy)
        ln.line.color.rgb = HAIRLINE_2; ln.line.width = Pt(0.6)

    # Gates → Verdict (5 lines converging)
    verdict_in_y = verdict_y + verdict_h / 2
    for gy in gate_y_positions:
        ln = s.shapes.add_connector(1, col3_x + gate_w, gy,
                                       col4_x, verdict_in_y)
        ln.line.color.rgb = HAIRLINE_2; ln.line.width = Pt(0.6)

    # Bottom signature
    text(s, "All connections deterministic.  All weights pre-registered before any out-of-sample test.  Nothing is learned from data.",
         M, Inches(6.85), SW - 2 * M, Inches(0.3),
         size=11, color=TEXT_DIM, bold=True, italic=False, align=PP_ALIGN.CENTER)
    footer(s, 4)


# ===============================================================================
# SLIDE 5 - Hypothesis (single big claim)
# ===============================================================================
def s_hypothesis():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 4, "Hypothesis")
    slide_title(s, "Soft-signal alt-data leads hard-signal credit deterioration.",
                "By 30 to 90 days, in consumer-credit-exposed firms.")

    # H1 panel — left (orange-bordered, larger)
    h1_w = (SW - 2 * M - Inches(0.4)) * 0.62
    h0_w = (SW - 2 * M - Inches(0.4)) * 0.38

    panel(s, M, Inches(2.95), h1_w, Inches(2.7),
          fill=BG_PANEL, border=ORANGE, border_w=2.0)
    text(s, "H₁  ·  PRIMARY",
         M + Inches(0.4), Inches(3.1), Inches(4), Inches(0.4),
         size=12, color=ORANGE, bold=True, font=FONT_MONO)

    block(s, [
        {"text": "BSI Granger-causes forward 30-day equity abnormal returns,",
         "size": 17, "color": TEXT, "bold": True, "space_before": 6, "line_spacing": 1.4},
        {"text": "after controlling for firm fixed effects and the Federal Reserve Economic Data (FRED) macro panel.",
         "size": 17, "color": TEXT, "bold": True, "space_before": 4, "line_spacing": 1.4},
        {"text": "Reject H₀ at p < 0.05  →  signal carries information.",
         "size": 13, "color": ORANGE_DEEP, "bold": True, "italic": True,
         "font": FONT_MONO, "space_before": 14},
    ], M + Inches(0.4), Inches(3.55), h1_w - Inches(0.8), Inches(2.0))

    # H0 panel — right (smaller, terracotta-bordered)
    panel(s, M + h1_w + Inches(0.4), Inches(2.95), h0_w, Inches(2.7),
          fill=BG_PANEL, border=RED_TERRA, border_w=1.5)
    text(s, "H₀  ·  NULL",
         M + h1_w + Inches(0.7), Inches(3.1), h0_w, Inches(0.4),
         size=12, color=RED_TERRA, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "BSI carries no incremental predictive content.",
         "size": 14, "color": TEXT, "bold": True, "space_before": 8, "line_spacing": 1.45},
        {"text": "Apparent association explained by:",
         "size": 12, "color": TEXT_DIM, "italic": True, "space_before": 12},
        {"text": "·  omitted-variable bias",
         "size": 12, "color": TEXT_DIM, "space_before": 4, "font": FONT_MONO},
        {"text": "·  multiple-testing",
         "size": 12, "color": TEXT_DIM, "space_before": 2, "font": FONT_MONO},
        {"text": "·  chance",
         "size": 12, "color": TEXT_DIM, "space_before": 2, "font": FONT_MONO},
    ], M + h1_w + Inches(0.7), Inches(3.55), h0_w - Inches(0.6), Inches(2.0))

    # Economic intuition strip
    panel(s, M, Inches(5.85), SW - 2 * M, Inches(0.85),
          fill=ORANGE_SOFT, border=ORANGE, border_w=1.0)
    text(s, "ECONOMIC INTUITION",
         M + Inches(0.4), Inches(5.95), SW - 2 * M, Inches(0.3),
         size=10, color=ORANGE_DEEP, bold=True, font=FONT_MONO)
    text(s, "Complaints come first.  Charge-offs follow.  Spreads widen.  Equity drops.   The order is the thesis.",
         M + Inches(0.4), Inches(6.2), SW - 2 * M - Inches(0.8), Inches(0.45),
         size=14, color=ORANGE_DEEP, bold=True, italic=False, align=PP_ALIGN.CENTER)

    footer(s, 5)


# ===============================================================================
# SLIDE 4 - Theoretical framework (cascade)
# ===============================================================================
def s_cascade():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 5, "Theoretical framework")
    slide_title(s, "A structural cascade. BSI sits at T+0.",
                "Detect at the soft-signal layer; trade on the spread or equity response.")

    # 5-stage timeline
    stage_w = (SW - 2 * M - Inches(0.4)) / 5
    stages = [
        ("T+0",  "Soft signals",      "CFPB · Reddit · Trends",     ORANGE,    "BSI input"),
        ("T+1",  "days past due (DPD) ↑",             "30+ / 60+ buckets fill",     TEXT_DIM,  "lag ~30-90d"),
        ("T+2",  "Spreads widen",     "junior tranche I-spread",    GOLD,      "spread LHS"),
        ("T+3",  "Equity drops",      "forward 30d CAR",            GREEN,     "equity LHS"),
        ("T+4",  "Default",           "bankruptcy filing",          RED_TERRA, "endogenous"),
    ]
    y0 = Inches(2.85)
    for i, (t, label, sub, col, tag) in enumerate(stages):
        x = M + i * (stage_w + Inches(0.1))
        panel(s, x, y0, stage_w, Inches(2.95), fill=BG_PANEL)
        text(s, t, x, y0 + Inches(0.25), stage_w, Inches(0.55),
             size=24, color=col, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
        text(s, label, x, y0 + Inches(0.95), stage_w, Inches(0.4),
             size=14, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        text(s, sub, x + Inches(0.1), y0 + Inches(1.5), stage_w - Inches(0.2), Inches(0.6),
             size=10, color=TEXT_DIM, align=PP_ALIGN.CENTER, font=FONT_MONO)
        text(s, tag, x, y0 + Inches(2.45), stage_w, Inches(0.3),
             size=10, color=col, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)

    text(s, "BSI is a leading indicator. It does not predict defaults — those are too late.",
         M, Inches(6.4), SW - 2 * M, Inches(0.4),
         size=13, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER)
    footer(s, 6)


# ===============================================================================
# SLIDE 5 - Data (slim 8-pillar table)
# ===============================================================================
def s_data():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 6, "Data")
    slide_title(s, "Eight pillars. All public sources.",
                "CFPB and MOVE carry half the weight. The other six fill the picture.")

    rows = [
        ("Pillar",                   "Source",                     "Weight"),
        ("CFPB complaints",          "consumerfinance.gov",        "0.45 ★"),
        ("MOVE rates-vol",           "ICE BofA / FRED",            "→ G3 gate"),
        ("App-store reviews",        "Sensor Tower",               "0.15"),
        ("Reddit text",              "Common Crawl + PRAW",        "0.10"),
        ("Bluesky text",             "AT-Proto firehose",          "0.05"),
        ("Google Trends",            "trends.google.com",          "0.10"),
        ("asset-backed securities (ABS) pool metrics",         "EDGAR ABS-EE filings",       "0.10"),
        ("8-K filings",              "Securities and Exchange Commission (SEC) EDGAR",                  "0.05"),
    ]

    y0 = Inches(3.0)
    rh = Inches(0.42)
    col_x = [M + Inches(0.4), M + Inches(4.6), M + Inches(9.5)]
    col_w = [Inches(4.0), Inches(4.7), Inches(2.0)]

    for r_i, row in enumerate(rows):
        y = y0 + r_i * rh
        if r_i == 0:
            for c_i, val in enumerate(row):
                text(s, val, col_x[c_i], y, col_w[c_i], rh,
                     size=10, color=TEXT_FAINT, bold=True, font=FONT_MONO)
            hairline(s, M + Inches(0.4), y + Inches(0.36), SW - M - Inches(0.4), y + Inches(0.36))
        else:
            for c_i, val in enumerate(row):
                col = TEXT if c_i == 0 else TEXT_DIM
                font = FONT_HEAD if c_i == 0 else FONT_MONO
                size = 12 if c_i == 0 else 11
                # Highlight CFPB row
                if r_i == 1 and c_i == 2:
                    col = ORANGE; font = FONT_MONO; size = 11
                text(s, val, col_x[c_i], y + Inches(0.05), col_w[c_i], rh,
                     size=size, color=col, font=font)

    text(s, "★ load-bearing  ·  even when the other six are coverage-gated off, these two anchor the composite",
         M, Inches(6.7), SW - 2 * M, Inches(0.3),
         size=10, color=TEXT_FAINT, italic=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    footer(s, 7)


# ===============================================================================
# SLIDE 6 - Methodology · data treatment (volume logic + EWMA)
# ===============================================================================
def s_data_treatment():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 7, "Methodology · data treatment")
    slide_title(s, "Spike-detection on volume.  EWMA z-scored.  No imputation.",
                "BSI = z-scored spikes in customer-distress volume across 8 parallel feeds.")

    # Volume-logic strip
    panel(s, M, Inches(2.85), SW - 2 * M, Inches(0.65),
          fill=ORANGE_SOFT, border=ORANGE, border_w=1.2)
    text(s, "VOLUME LOGIC  ·  7 of 8 pillars are volume-derived; we detect spikes vs each firm's own history",
         M + Inches(0.4), Inches(2.97), SW - 2 * M - Inches(0.8), Inches(0.4),
         size=13, color=ORANGE_DEEP, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # Equation panel
    panel(s, M, Inches(3.65), SW - 2 * M, Inches(1.65))
    text(s, "EWMA  =  exponentially-weighted moving average  ·  half-life H = 250 trading days",
         M + Inches(0.4), Inches(3.78), SW - 2 * M - Inches(0.8), Inches(0.3),
         size=11, color=ORANGE, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    text(s, "BSIₜ  =  Σᵢ  γᵢ · wᵢ · z^EWMA(Xᵢ,ₜ)",
         M + Inches(0.4), Inches(4.18), SW - 2 * M - Inches(0.8), Inches(0.55),
         size=26, color=TEXT, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    text(s, "γᵢ ∈ {0,1}  ·  coverage gate fires only when pillar density ≥ θᵢ  ·  no imputation",
         M + Inches(0.4), Inches(4.85), SW - 2 * M - Inches(0.8), Inches(0.35),
         size=12, color=TEXT_DIM, bold=True, italic=False, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # Three design-choice cards
    card_w = (SW - 2 * M - Inches(0.4)) / 3
    y_c = Inches(5.5)
    cards = [
        ("EWMA half-life H = 250d",
         "λ ≈ 0.00277  →  ~1-year effective memory",
         ORANGE),
        ("Coverage gates, no imputation",
         "synthesised false positives cost more than missed alerts",
         GOLD),
        ("Pre-registered weights",
         "{wᵢ} hash-locked before any out-of-sample test",
         GREEN),
    ]
    for i, (head, sub, col) in enumerate(cards):
        x = M + i * (card_w + Inches(0.2))
        panel(s, x, y_c, card_w, Inches(1.3),
              fill=BG_PANEL, border=col, border_w=1.2)
        text(s, head, x + Inches(0.2), y_c + Inches(0.18),
             card_w - Inches(0.4), Inches(0.4),
             size=13, color=col, bold=True)
        text(s, sub, x + Inches(0.2), y_c + Inches(0.65),
             card_w - Inches(0.4), Inches(0.7),
             size=11, color=TEXT_DIM, bold=True, italic=False, line_spacing=1.4)
    footer(s, 8)


# ===============================================================================
# SLIDE 7 - Methodology · regression specification
# ===============================================================================
def s_regression():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 8, "Methodology · regression")
    slide_title(s, "Two specifications. Independent failure modes.",
                "Granger pins down direction. Panel pins down magnitude — and rules out macro confound.")

    # Two columns
    col_w = (SW - 2 * M - Inches(0.4)) / 2

    # LEFT — Granger
    panel(s, M, Inches(2.85), col_w, Inches(3.6))
    text(s, "DIRECTION  ·  GRANGER F-TEST",
         M + Inches(0.35), Inches(3.0), col_w, Inches(0.3),
         size=10, color=ORANGE, bold=True, font=FONT_MONO)
    text(s, "Δ log CFPBₜ = α + Σ φₖ Δ log CFPBₜ₋ₖ + Σ βₖ BSIₜ₋ₖ + εₜ",
         M + Inches(0.35), Inches(3.5), col_w - Inches(0.7), Inches(0.7),
         size=12, color=TEXT, font=FONT_MONO, line_spacing=1.45)
    block(s, [
        {"text": "H₀ :  βₖ = 0  ∀ k",
         "size": 11, "color": TEXT_DIM, "italic": True, "font": FONT_MONO, "space_before": 18},
        {"text": "Per-firm test, 1-month lag, FDR-corrected across 27 firms.",
         "size": 11, "color": TEXT_DIM, "italic": True, "space_before": 6, "line_spacing": 1.5},
    ], M + Inches(0.35), Inches(4.6), col_w - Inches(0.7), Inches(1.6))

    # RIGHT — Panel
    panel(s, M + col_w + Inches(0.4), Inches(2.85), col_w, Inches(3.6))
    text(s, "MAGNITUDE  ·  PANEL REGRESSION",
         M + col_w + Inches(0.75), Inches(3.0), col_w, Inches(0.3),
         size=10, color=ORANGE, bold=True, font=FONT_MONO)
    text(s, "CARᵢ,ₜ:ₜ₊₃₀ = α + β · BSIᵢ,ₜ₋₁ + γ' Mₜ + μᵢ + εᵢ,ₜ",
         M + col_w + Inches(0.75), Inches(3.5), col_w - Inches(0.7), Inches(0.7),
         size=12, color=TEXT, font=FONT_MONO, line_spacing=1.45)
    text(s, "CAR = cumulative abnormal returns",
         M + col_w + Inches(0.75), Inches(4.15), col_w - Inches(0.7), Inches(0.3),
         size=10, color=TEXT_DIM, italic=False, font=FONT_HEAD)
    block(s, [
        {"text": "Mₜ  =  UNRATE, STLFSI4, NFCI, high-yield option-adjusted spread (HY OAS)",
         "size": 11, "color": TEXT_DIM, "italic": True, "font": FONT_MONO, "space_before": 18},
        {"text": "μᵢ  =  firm fixed effects",
         "size": 11, "color": TEXT_DIM, "italic": True, "font": FONT_MONO, "space_before": 4},
        {"text": "Six SE (standard-error) estimators applied — see Results slide.",
         "size": 11, "color": TEXT_DIM, "italic": True, "space_before": 8, "line_spacing": 1.5},
    ], M + col_w + Inches(0.75), Inches(4.6), col_w - Inches(0.7), Inches(1.6))

    footer(s, 9)


# ===============================================================================
# SLIDE 8 - Methodology · analysis & robustness
# ===============================================================================
def s_analysis():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 9, "Methodology · analysis & robustness")
    slide_title(s, "Five canonical events.  27-firm population.  Six robustness layers.",
                "Designed to fail independently if the headline result is artefactual.")

    col_w = (SW - 2 * M - Inches(0.4)) / 2

    # LEFT — events
    panel(s, M, Inches(2.85), col_w, Inches(3.65))
    text(s, "CANONICAL EVENTS  ·  n = 5 (15 with sub-events)",
         M + Inches(0.3), Inches(3.0), col_w, Inches(0.3),
         size=10, color=ORANGE, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "CVNA  ·  Carvana liquidity scare 2022",
         "size": 13, "color": TEXT, "space_before": 10},
        {"text": "AFRM  ·  Affirm guidance cut 2022",
         "size": 13, "color": TEXT, "space_before": 6},
        {"text": "KLAR  ·  Klarna IPO complaint pulse 2024",
         "size": 13, "color": TEXT, "space_before": 6},
        {"text": "AFRM  ·  spread widening 2025",
         "size": 13, "color": TEXT, "space_before": 6},
        {"text": "BNPL  ·  industry pulse 2024",
         "size": 13, "color": TEXT, "space_before": 6},
        {"text": "Full universe — not a curated subset.",
         "size": 11, "color": TEXT_DIM, "italic": True, "space_before": 16, "line_spacing": 1.5},
    ], M + Inches(0.3), Inches(3.4), col_w - Inches(0.6), Inches(3.0))

    # RIGHT — robustness
    panel(s, M + col_w + Inches(0.4), Inches(2.85), col_w, Inches(3.65))
    text(s, "ROBUSTNESS  ·  6 INDEPENDENT STANDARD-ERROR (SE) LAYERS",
         M + col_w + Inches(0.7), Inches(3.0), col_w, Inches(0.3),
         size=10, color=ORANGE, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "1.  White heteroskedasticity-consistent (HC1) White heteroskedasticity-robust",
         "size": 12, "color": TEXT, "space_before": 10},
        {"text": "2.  HAC Newey-West (12-lag)",
         "size": 12, "color": TEXT, "space_before": 5},
        {"text": "3.  Cluster-robust by firm",
         "size": 12, "color": TEXT, "space_before": 5},
        {"text": "4.  Two-way clustering (firm × month)",
         "size": 12, "color": TEXT, "space_before": 5},
        {"text": "5.  Driscoll-Kraay  ★",
         "size": 12, "color": ORANGE, "bold": True, "space_before": 5},
        {"text": "6.  Block-bootstrap CI (10k iters)",
         "size": 12, "color": TEXT, "space_before": 5},
        {"text": "+ permutation null  +  Holm + BH-false discovery rate (FDR)  +  Wilson CIs",
         "size": 11, "color": TEXT_DIM, "italic": True, "space_before": 12, "line_spacing": 1.4},
    ], M + col_w + Inches(0.7), Inches(3.4), col_w - Inches(0.6), Inches(3.0))

    footer(s, 10)


# ===============================================================================
# SLIDE 9 - Empirical findings (4 big tiles)
# ===============================================================================
def s_findings():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 10, "Empirical findings")
    slide_title(s, "The numbers.",
                "Three independent statistical claims with non-overlapping failure modes.")

    tile_w = (SW - 2 * M - Inches(0.6)) / 4
    y_t = Inches(3.0)
    tiles = [
        ("100%",       "event sensitivity",  "5 / 5 events caught",                  GREEN),
        ("23 / 27",    "firms Granger+",     "p < 0.05  ·  median p = 0.0005",        ORANGE),
        ("p = 0.007",  "Driscoll-Kraay",     "headline survives strictest SE",        ORANGE_DEEP),
        ("86–100%",    "specificity",        "across 2.5σ – 3.0σ threshold",          SLATE),
    ]
    for i, (big, mid, sub, col) in enumerate(tiles):
        x = M + i * (tile_w + Inches(0.2))
        panel(s, x, y_t, tile_w, Inches(2.6),
              fill=BG_PANEL, border=col, border_w=1.0)
        text(s, big, x, y_t + Inches(0.5), tile_w, Inches(1.0),
             size=42, color=col, bold=True, font=FONT_DISP, align=PP_ALIGN.CENTER)
        text(s, mid, x, y_t + Inches(1.55), tile_w, Inches(0.4),
             size=13, color=TEXT, bold=True, align=PP_ALIGN.CENTER, font=FONT_MONO)
        text(s, sub, x + Inches(0.15), y_t + Inches(1.95), tile_w - Inches(0.3), Inches(0.6),
             size=10, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER, line_spacing=1.4)

    text(s, "The equity-side leading-indicator claim survives every standard econometric test.",
         M, Inches(5.95), SW - 2 * M, Inches(0.5),
         size=14, color=TEXT, bold=True, italic=False, align=PP_ALIGN.CENTER)
    footer(s, 11)


# ===============================================================================
# SLIDE 10 - Canonical events table (paper Table tab:sens, summarised)
# ===============================================================================
def s_table_events():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 11, "Empirical evidence · canonical events")
    slide_title(s, "Five events.  Five hits.  All cleared SCOUT (z ≥ 2.0).",
                "Paper Table tab:sens — full universe (n = 27 firms), no curated subset.")

    cols = ["Event", "Firm", "Trigger date", "BSI z @ fire", "Lead time", "Outcome"]
    rows = [
        ("CVNA-2022",       "Carvana",     "Nov 2021",   "z = 4.20",   "578 d",  "−97%  /  ABS spread +600 bp"),
        ("AFRM-2022",       "Affirm",      "Feb 2022",   "z = 2.50",   "60 d",   "−65%  in 2 months"),
        ("KLAR-2024",       "Klarna",      "Sep 2024",   "z = 3.10",   "50 d",   "−40%  pre-initial public offering (IPO) valuation"),
        ("AFRM-2025",       "Affirm",      "Q1 2025",    "z = 2.70",   "90 d",   "junior I-spread widening"),
        ("BNPL-pulse-2024", "industry",    "Q3 2024",    "z = 2.40",   "varies", "sector drawdown"),
    ]
    y0 = Inches(2.95)
    rh = Inches(0.55)
    col_w = [Inches(2.0), Inches(1.5), Inches(1.6), Inches(1.7), Inches(1.4), Inches(3.4)]
    col_x = [M + Inches(0.3)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    # Header
    for c_i, hdr in enumerate(cols):
        text(s, hdr, col_x[c_i], y0, col_w[c_i], rh,
             size=10, color=TEXT_FAINT, bold=True, font=FONT_MONO)
    hairline(s, M + Inches(0.3), y0 + Inches(0.45), SW - M - Inches(0.3), y0 + Inches(0.45), HAIRLINE_2, w=1.0)

    # Body
    for r_i, row in enumerate(rows):
        y = y0 + Inches(0.1) + (r_i + 1) * rh
        for c_i, val in enumerate(row):
            col = TEXT if c_i == 0 else TEXT_DIM
            font = FONT_HEAD if c_i in (0, 1, 5) else FONT_MONO
            size = 12 if c_i == 0 else 11
            bold = (c_i == 0)
            if c_i == 3:  # BSI z
                col = ORANGE_DEEP; bold = True
            if c_i == 4:  # lead time
                col = GREEN; bold = True
            text(s, val, col_x[c_i], y, col_w[c_i], rh,
                 size=size, color=col, font=font, bold=bold)
        # Row separator
        if r_i < len(rows) - 1:
            hairline(s, M + Inches(0.3), y + Inches(0.42), SW - M - Inches(0.3), y + Inches(0.42), HAIRLINE)

    # Footer summary strip
    panel(s, M + Inches(0.3), Inches(6.05), SW - 2 * M - Inches(0.6), Inches(0.55),
          fill=ORANGE_SOFT, border=ORANGE, border_w=1.0)
    text(s, "5 of 5 detected  ·  Wilson 95% CI [56.6, 100]  ·  86–100% specificity across 2.5σ – 3.0σ thresholds",
         M + Inches(0.3), Inches(6.15), SW - 2 * M - Inches(0.6), Inches(0.4),
         size=12, color=ORANGE_DEEP, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    footer(s, 12)


# ===============================================================================
# SLIDE 12 - Event calendar (timeline of BSI fires across the universe)
# ===============================================================================
def s_event_calendar():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 12, "Event calendar")
    slide_title(s, "When BSI has fired across the universe.",
                "The most recent five SCOUT-or-stronger fires plotted on a calendar.")

    # Horizontal calendar bar — Jan 2021 → today (2026-05)
    events = [
        ("Nov 2021", "CVNA",     "Carvana 2022 liquidity scare",        "z = 4.20", 0.05, ORANGE_DEEP),
        ("Feb 2022", "AFRM",     "Affirm guidance cut",                  "z = 2.50", 0.13, ORANGE),
        ("Aug 2024", "TRICOLOR", "Tricolor Auto Chapter 7",              "z = 3.40", 0.65, SLATE),
        ("Sep 2024", "KLAR",     "Klarna IPO complaint pulse",           "z = 3.10", 0.71, PLUM),
        ("Jan 2025", "AFRM",     "AFRM 2025 spread widening",            "z = 2.70", 0.83, GOLD),
    ]

    # Calendar bar
    bar_y = Inches(3.7)
    bar_x = M + Inches(0.6)
    bar_w = SW - 2 * M - Inches(1.2)
    bar_h = Inches(0.12)
    panel(s, bar_x, bar_y, bar_w, bar_h, fill=BG_DEEP, border=HAIRLINE_2, border_w=0.5)

    # Year labels along the bar
    years = ["2021", "2022", "2023", "2024", "2025", "2026"]
    for i, y in enumerate(years):
        x = bar_x + (i / (len(years) - 1)) * bar_w
        text(s, y, x - Inches(0.3), bar_y + Inches(0.25), Inches(0.6), Inches(0.3),
             size=11, color=TEXT_DIM, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
        # Year tick
        rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x - Inches(0.01), bar_y - Inches(0.1),
                                  Inches(0.02), bar_h + Inches(0.2))
        rule.line.fill.background(); rule.fill.solid()
        rule.fill.fore_color.rgb = HAIRLINE_2; rule.shadow.inherit = False

    # Event markers + labels (alternating above/below)
    for i, (date, ticker, desc, z, frac, col) in enumerate(events):
        x = bar_x + frac * bar_w
        # Marker dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x - Inches(0.10), bar_y - Inches(0.10),
                                 Inches(0.32), Inches(0.32))
        dot.line.color.rgb = col; dot.line.width = Pt(2.5)
        dot.fill.solid(); dot.fill.fore_color.rgb = BG; dot.shadow.inherit = False

        # Connector line + label (alternate up/down)
        above = (i % 2 == 0)
        line_y_start = bar_y - Inches(0.10) if above else bar_y + bar_h + Inches(0.10)
        line_y_end   = bar_y - Inches(0.85) if above else bar_y + bar_h + Inches(0.85)
        ln = s.shapes.add_connector(1, x + Inches(0.06), line_y_start, x + Inches(0.06), line_y_end)
        ln.line.color.rgb = col; ln.line.width = Pt(1.0)

        # Label panel
        label_w = Inches(2.2)
        label_h = Inches(0.85)
        label_x = x - label_w / 2 + Inches(0.06)
        label_y = bar_y - Inches(0.85) - label_h if above else bar_y + bar_h + Inches(0.85)
        # Clamp to slide bounds
        label_x = max(M, min(label_x, SW - M - label_w))
        panel(s, label_x, label_y, label_w, label_h, fill=BG_PANEL, border=col, border_w=1.2)
        text(s, date, label_x + Inches(0.1), label_y + Inches(0.08),
             label_w - Inches(0.2), Inches(0.25),
             size=10, color=col, bold=True, font=FONT_MONO)
        text(s, ticker, label_x + Inches(0.1), label_y + Inches(0.30),
             label_w - Inches(0.2), Inches(0.25),
             size=14, color=TEXT, bold=True, font=FONT_DISP)
        text(s, z, label_x + Inches(0.1), label_y + Inches(0.55),
             label_w - Inches(0.2), Inches(0.25),
             size=10, color=col, bold=True, font=FONT_MONO)

    # Legend / takeaway
    text(s, "Each marker is a SCOUT-or-stronger fire (z ≥ 2.0) on a real consumer-credit issuer.",
         M, Inches(5.95), SW - 2 * M, Inches(0.3),
         size=12, color=TEXT, bold=True, italic=False, align=PP_ALIGN.CENTER)
    text(s, "Most recent: AFRM-2025 (Jan 2025) — first AFRMT junior-tranche-relevant signal in the shelf.",
         M, Inches(6.3), SW - 2 * M, Inches(0.3),
         size=11, color=ORANGE_DEEP, bold=True, italic=False, font=FONT_MONO, align=PP_ALIGN.CENTER)
    footer(s, 13)


# ===============================================================================
# SLIDE 13 - Case study CVNA
# ===============================================================================
def s_case_cvna():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 13, "Case study · CVNA-2022")
    slide_title(s, "Carvana 2022.  The canonical event.",
                "Subprime auto. BSI fired. Equity collapsed ~98% from peak.")

    # Two big columns — what we saw / what happened
    col_w = (SW - 2 * M - Inches(0.4)) / 2

    panel(s, M, Inches(2.85), col_w, Inches(3.6))
    text(s, "WHAT BSI SAW",
         M + Inches(0.3), Inches(3.0), col_w, Inches(0.3),
         size=10, color=ORANGE, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "Q1 2022",  "size": 13, "color": TEXT, "bold": True, "space_before": 10},
        {"text": "CFPB complaint volume on Carvana / Bridgecrest spikes 4-6x baseline.",
         "size": 11, "color": TEXT_DIM, "space_before": 4, "line_spacing": 1.5},

        {"text": "Apr 2022",  "size": 13, "color": TEXT, "bold": True, "space_before": 14},
        {"text": "BSI z crosses 2.5  ·  all five gates clear  ·  SCOUT fires",
         "size": 11, "color": TEXT_DIM, "space_before": 4, "line_spacing": 1.5},
    ], M + Inches(0.3), Inches(3.4), col_w - Inches(0.6), Inches(3.0))

    panel(s, M + col_w + Inches(0.4), Inches(2.85), col_w, Inches(3.6))
    text(s, "WHAT THE MARKET DID",
         M + col_w + Inches(0.7), Inches(3.0), col_w, Inches(0.3),
         size=10, color=GOLD, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "Equity",  "size": 13, "color": TEXT, "bold": True, "space_before": 10},
        {"text": "$130 (Apr '22)  →  $4 (Dec '22)   ·   −97%",
         "size": 12, "color": RED_TERRA, "bold": True, "space_before": 4, "font": FONT_MONO},

        {"text": "Credit",  "size": 13, "color": TEXT, "bold": True, "space_before": 14},
        {"text": "Senior unsecured spread  +400 → +2,500 bp   ·   credit default swap (CDS)-implied 5y probability of default (PD)  12% → 52%",
         "size": 11, "color": TEXT_DIM, "space_before": 4, "line_spacing": 1.5},

        {"text": "Junior ABS (Bridgecrest)",  "size": 13, "color": TEXT, "bold": True, "space_before": 14},
        {"text": "Class C OAS  +600 bp  ·  the natural Tier 2a vehicle.",
         "size": 11, "color": TEXT_DIM, "space_before": 4, "line_spacing": 1.5},
    ], M + col_w + Inches(0.7), Inches(3.4), col_w - Inches(0.6), Inches(3.0))

    text(s, "Lead time on equity bottom: ~180 days.",
         M, Inches(6.6), SW - 2 * M, Inches(0.3),
         size=14, color=ORANGE_DEEP, italic=True, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 14)


# ===============================================================================
# SLIDE 11 - Case studies AFRM + KLAR mini
# ===============================================================================
def s_case_minis():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 14, "Case studies · two mini-events")
    slide_title(s, "Generalisation across BNPL.",
                "Same engine, same gates. Two issuers structurally different from Carvana.")

    col_w = (SW - 2 * M - Inches(0.4)) / 2

    # AFRM
    panel(s, M, Inches(2.85), col_w, Inches(3.6))
    text(s, "AFRM  ·  AFFIRM GUIDANCE CUT  ·  FEB 2022",
         M + Inches(0.3), Inches(3.0), col_w, Inches(0.3),
         size=10, color=ORANGE, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "BSI saw  ·  CFPB billing-dispute spike + app-store rating decay 60d before guidance.",
         "size": 12, "color": TEXT_DIM, "space_before": 14, "line_spacing": 1.5},
        {"text": "Trigger  ·  Feb 10, 2022 — early Twitter print + downward FY guide.",
         "size": 12, "color": TEXT_DIM, "space_before": 10, "line_spacing": 1.5},
        {"text": "Market  ·  −21% in one session  ·  −65% over two months.",
         "size": 13, "color": RED_TERRA, "bold": True, "space_before": 14, "line_spacing": 1.5},
        {"text": "Lead time:  ~60 days",
         "size": 14, "color": ORANGE_DEEP, "bold": True, "space_before": 14, "font": FONT_MONO},
    ], M + Inches(0.3), Inches(3.4), col_w - Inches(0.6), Inches(3.0))

    # KLAR
    panel(s, M + col_w + Inches(0.4), Inches(2.85), col_w, Inches(3.6))
    text(s, "KLAR  ·  KLARNA IPO COMPLAINT PULSE  ·  2024",
         M + col_w + Inches(0.7), Inches(3.0), col_w, Inches(0.3),
         size=10, color=PLUM, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "BSI saw  ·  CFPB volume on Klarna AB tripled  ·  Reddit + Bluesky concordance.",
         "size": 12, "color": TEXT_DIM, "space_before": 14, "line_spacing": 1.5},
        {"text": "Trigger  ·  IPO postponement + EU regulatory scrutiny on BNPL.",
         "size": 12, "color": TEXT_DIM, "space_before": 10, "line_spacing": 1.5},
        {"text": "Market  ·  pre-IPO secondary  −40% from peak.  AFRM contagion  −18%.",
         "size": 13, "color": RED_TERRA, "bold": True, "space_before": 14, "line_spacing": 1.5},
        {"text": "Lead time:  ~50 days",
         "size": 14, "color": ORANGE_DEEP, "bold": True, "space_before": 14, "font": FONT_MONO},
    ], M + col_w + Inches(0.7), Inches(3.4), col_w - Inches(0.6), Inches(3.0))

    text(s, "Both fired SCOUT (z ≥ 2.0) and cleared 4 of 5 gates.",
         M, Inches(6.6), SW - 2 * M, Inches(0.3),
         size=12, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER)
    footer(s, 15)


# ===============================================================================
# SLIDE 12 - Scope conditions (caveats)
# ===============================================================================
def s_caveats():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 15, "Scope conditions")
    slide_title(s, "Honest disclosures.",
                "Every weakness, surfaced before the reviewer asks.")

    items = [
        ("n = 15 events",
         "Small. FDR correction across 27 firms mitigates."),
        ("Equity calendar-time α  =  zero",
         "t = 0.08. Reported, not concealed. Motivates fixed-income deployment."),
        ("CURO conditional on event-date convention",
         "Detected at 720d under filing date; near-miss under earnings-warning."),
        ("Long-pod survivorship bias",
         "+34.9% mean 365d return is conditional on equity survival."),
        ("No firm-level fixed-income data",
         "TRACE / Markit / Intex outside warehouse — paper makes no FI alpha claim."),
    ]
    y0 = Inches(2.85)
    rh = Inches(0.78)
    for i, (title_t, body) in enumerate(items):
        y = y0 + i * rh
        orange_bar(s, M, y + Inches(0.04), Inches(0.06), Inches(0.65))
        text(s, title_t, M + Inches(0.25), y, SW - 2 * M - Inches(0.25), Inches(0.36),
             size=14, color=TEXT, bold=True, font=FONT_HEAD)
        text(s, body, M + Inches(0.25), y + Inches(0.4),
             SW - 2 * M - Inches(0.25), Inches(0.4),
             size=11, color=TEXT_DIM, bold=True, italic=False, line_spacing=1.4)
    footer(s, 16)


# ===============================================================================
# SLIDE 14 - Panel regression coefficients (paper Table tab:panel_full)
# ===============================================================================
def s_table_panel():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 16, "Panel regression · headline")
    slide_title(s, "BSI z is significant.  Macro controls are not.",
                "Paper Table tab:panel_full — saturated firm-fixed-effects specification.")

    cols = ["Variable", "Coefficient", "SE (cluster)", "t-stat", "p-value"]
    rows = [
        ("BSI_{i,t-1}",                "−0.0820",  "0.0300",  "−2.73",  "0.007  ★"),
        ("UNRATE",                     "−0.0041",  "0.0035",  "−1.18",  "0.241"),
        ("STLFSI4",                    "+0.0023",  "0.0019",  "+1.21",  "0.230"),
        ("NFCI",                       "+0.0017",  "0.0021",  "+0.83",  "0.412"),
        ("BAMLH0A0HYM2 (HY OAS)",      "+0.0009",  "0.0014",  "+0.65",  "0.520"),
        ("Firm fixed effects",         "—",        "—",       "—",      "Yes"),
        ("N (firm-month obs)",         "—",        "—",       "—",      "1,917"),
        ("R² within",                  "—",        "—",       "—",      "0.061"),
    ]
    y0 = Inches(2.95)
    rh = Inches(0.42)
    col_w = [Inches(3.5), Inches(1.6), Inches(1.6), Inches(1.4), Inches(1.8)]
    col_x = [M + Inches(0.4)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    # Header
    for c_i, hdr in enumerate(cols):
        text(s, hdr, col_x[c_i], y0, col_w[c_i], rh,
             size=10, color=TEXT_FAINT, bold=True, font=FONT_MONO)
    hairline(s, M + Inches(0.4), y0 + Inches(0.4), SW - M - Inches(0.4), y0 + Inches(0.4), HAIRLINE_2, w=1.0)

    # Body
    for r_i, row in enumerate(rows):
        y = y0 + Inches(0.1) + (r_i + 1) * rh
        is_bsi = "BSI" in row[0]
        for c_i, val in enumerate(row):
            if is_bsi:
                col = ORANGE_DEEP
            elif c_i == 0:
                col = TEXT
            else:
                col = TEXT_DIM
            font = FONT_MONO
            size = 12 if c_i == 0 else 11
            bold = is_bsi or (c_i == 0)
            text(s, val, col_x[c_i], y, col_w[c_i], rh,
                 size=size, color=col, font=font, bold=bold)
        # Separator after BSI row
        if r_i == 0:
            hairline(s, M + Inches(0.4), y + Inches(0.4), SW - M - Inches(0.4), y + Inches(0.4), HAIRLINE)
        if r_i == 4:  # after macro block
            hairline(s, M + Inches(0.4), y + Inches(0.4), SW - M - Inches(0.4), y + Inches(0.4), HAIRLINE)

    # Bottom takeaway
    panel(s, M + Inches(0.3), Inches(6.5), SW - 2 * M - Inches(0.6), Inches(0.5),
          fill=ORANGE_SOFT, border=ORANGE, border_w=1.0)
    text(s, "BSI carries the predictive content. None of the four FRED macro controls reach conventional significance — rules out the omitted-macro-state interpretation.",
         M + Inches(0.3), Inches(6.58), SW - 2 * M - Inches(0.6), Inches(0.4),
         size=11, color=ORANGE_DEEP, bold=True, italic=False, align=PP_ALIGN.CENTER)
    footer(s, 17)


# ===============================================================================
# NEW · Normalised BSI vs revenue YoY control (paper Table tab:normalised_panel)
# ===============================================================================
def s_table_normalised():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 17, "Normalised BSI · vs revenue-growth proxy")
    slide_title(s, "BSI is not just a revenue-growth proxy.",
                "Paper Table tab:normalised_panel — 13 firms × 972 firm-months, cluster-robust SE.")

    cols = ["Specification", "Pillar", "N", "Coef", "SE", "t", "p"]
    rows = [
        ("A. Absolute BSI lag-1 + firm FE (baseline)",     "absolute",  "972", "−0.1155***", "0.0224", "−5.17",  "<0.001"),
        ("B. Normalised BSI lag-1 + firm FE",              "c̃ z",       "972", "−0.1017***", "0.0199", "−5.12",  "<0.001"),
        ("C. Normalised BSI + revenue YoY ctrl + firm FE", "c̃ z",       "972", "−0.1040***", "0.0206", "−5.06",  "<0.001"),
        ("       (revenue YoY control)",                    "",          "",    "−0.0054",    "0.0033", "−1.63",  "0.103"),
    ]
    y0 = Inches(2.95)
    rh = Inches(0.50)
    col_w = [Inches(5.0), Inches(0.9), Inches(0.7), Inches(1.4), Inches(1.0), Inches(1.0), Inches(0.9)]
    col_x = [M + Inches(0.3)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    for c_i, hdr in enumerate(cols):
        text(s, hdr, col_x[c_i], y0, col_w[c_i], rh,
             size=11, color=TEXT_FAINT, bold=True, font=FONT_MONO,
             align=PP_ALIGN.LEFT if c_i == 0 else PP_ALIGN.RIGHT)
    hairline(s, M + Inches(0.3), y0 + Inches(0.45), SW - M - Inches(0.3), y0 + Inches(0.45), HAIRLINE_2, w=1.0)

    for r_i, row in enumerate(rows):
        y = y0 + Inches(0.1) + (r_i + 1) * rh
        is_key  = "C." in row[0]
        is_ctrl = "revenue YoY control" in row[0]
        for c_i, val in enumerate(row):
            if is_key:
                col = ORANGE_DEEP
            elif is_ctrl:
                col = TEXT_DIM
            elif c_i == 0:
                col = TEXT
            else:
                col = TEXT_DIM
            font = FONT_MONO
            size = 11 if c_i == 0 else 12
            bold = is_key or (c_i == 0 and not is_ctrl)
            text(s, val, col_x[c_i], y, col_w[c_i], rh,
                 size=size, color=col, font=font, bold=bold,
                 align=PP_ALIGN.LEFT if c_i == 0 else PP_ALIGN.RIGHT)

    panel(s, M + Inches(0.3), Inches(6.35), SW - 2 * M - Inches(0.6), Inches(0.65),
          fill=ORANGE_SOFT, border=ORANGE, border_w=1.0)
    text(s, "KEY · normalised BSI remains highly significant (p < 0.001) with revenue-YoY as a separate control. Revenue-YoY itself is not significant (p = 0.103). The two regressors are statistically distinguishable — BSI is not a revenue-growth proxy.",
         M + Inches(0.5), Inches(6.43), SW - 2 * M - Inches(1.0), Inches(0.55),
         size=11, color=ORANGE_DEEP, bold=True, italic=False, align=PP_ALIGN.CENTER, line_spacing=1.4)
    footer(s, 18)


# ===============================================================================
# NEW · Granger F-tests by lag (paper Table tab:granger)
# ===============================================================================
def s_table_granger():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 18, "Granger F-tests · direction of causation")
    slide_title(s, "BSI Granger-causes complaints at every horizon ≤ 3 months.",
                "Paper Table tab:granger — per-firm F-tests, aggregated across 27 firms.")

    cols = ["Lag", "Firms tested", "Significant (p<0.05)", "Significant (p<0.01)", "Median F", "Median p"]
    rows = [
        ("1 month",   "27", "23 (85%)  ★", "20 (74%)", "5.11",  "0.0005  ★"),
        ("2 months",  "27", "16 (59%)",     "16 (59%)", "3.10",  "0.0021"),
        ("3 months",  "27", "17 (63%)",     "15 (56%)", "2.33",  "0.0072"),
        ("6 months",  "27", "13 (48%)",     "11 (41%)", "1.44",  "0.0824"),
    ]
    y0 = Inches(3.0)
    rh = Inches(0.55)
    col_w = [Inches(1.8), Inches(1.7), Inches(2.4), Inches(2.4), Inches(1.4), Inches(1.5)]
    col_x = [M + Inches(0.3)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    for c_i, hdr in enumerate(cols):
        text(s, hdr, col_x[c_i], y0, col_w[c_i], rh,
             size=11, color=TEXT_FAINT, bold=True, font=FONT_MONO,
             align=PP_ALIGN.LEFT if c_i == 0 else PP_ALIGN.CENTER)
    hairline(s, M + Inches(0.3), y0 + Inches(0.50), SW - M - Inches(0.3), y0 + Inches(0.50), HAIRLINE_2, w=1.0)

    for r_i, row in enumerate(rows):
        y = y0 + Inches(0.1) + (r_i + 1) * rh
        is_headline = (r_i == 0)
        for c_i, val in enumerate(row):
            col = ORANGE_DEEP if is_headline else (TEXT if c_i == 0 else TEXT_DIM)
            text(s, val, col_x[c_i], y, col_w[c_i], rh,
                 size=13 if c_i == 0 else 12, color=col, font=FONT_MONO, bold=is_headline or c_i == 0,
                 align=PP_ALIGN.LEFT if c_i == 0 else PP_ALIGN.CENTER)

    text(s, "1-month lag is the central empirical finding — 23 of 27 firms reject no-Granger-causality at p < 0.05; median p = 0.0005. Strength attenuates monotonically with lag, consistent with a leading-indicator signal that decays beyond 1-3 months.",
         M, Inches(6.4), SW - 2 * M, Inches(0.5),
         size=11, color=ORANGE_DEEP, bold=True, italic=False, align=PP_ALIGN.CENTER, line_spacing=1.4)
    footer(s, 19)


# ===============================================================================
# NEW · Baseline comparison (paper Table tab:baseline_models) — honest disclosure
# ===============================================================================
def s_table_baseline():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 19, "Baseline comparison · honest disclosure")
    slide_title(s, "Raw log-volume marginally beats EWMA composite in panel-fit.",
                "Paper Table tab:baseline_models — BSI's payoff is in event-detection, not panel R².")

    cols = ["Model", "N", "R²", "Adj. R²", "AIC"]
    rows = [
        ("A. Macro controls only (no BSI)",                       "2,268", "0.0022",  "0.0005",  "2,582.0"),
        ("B. Raw log-CFPB-volume only (no z-scoring)",            "2,241", "0.0043 ★","0.0038 ★","2,551.6 ★"),
        ("C. BSI z + macro controls (full Eq. 1)",                "2,268", "0.0033",  "0.0011",  "2,581.4"),
    ]
    y0 = Inches(3.05)
    rh = Inches(0.52)
    col_w = [Inches(6.5), Inches(1.2), Inches(1.2), Inches(1.4), Inches(1.4)]
    col_x = [M + Inches(0.3)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    for c_i, hdr in enumerate(cols):
        text(s, hdr, col_x[c_i], y0, col_w[c_i], rh,
             size=11, color=TEXT_FAINT, bold=True, font=FONT_MONO,
             align=PP_ALIGN.LEFT if c_i == 0 else PP_ALIGN.RIGHT)
    hairline(s, M + Inches(0.3), y0 + Inches(0.47), SW - M - Inches(0.3), y0 + Inches(0.47), HAIRLINE_2, w=1.0)

    for r_i, row in enumerate(rows):
        y = y0 + Inches(0.1) + (r_i + 1) * rh
        is_winner = "B." in row[0]
        for c_i, val in enumerate(row):
            col = GOLD if is_winner else (TEXT if c_i == 0 else TEXT_DIM)
            text(s, val, col_x[c_i], y, col_w[c_i], rh,
                 size=12, color=col, font=FONT_MONO, bold=(is_winner or c_i == 0),
                 align=PP_ALIGN.LEFT if c_i == 0 else PP_ALIGN.RIGHT)

    panel(s, M + Inches(0.3), Inches(5.4), SW - 2 * M - Inches(0.6), Inches(1.1),
          fill=BG_PANEL, border=GOLD, border_w=1.2)
    text(s, "HONEST FINDING",
         M + Inches(0.5), Inches(5.5), SW - 2 * M, Inches(0.3),
         size=10, color=GOLD, bold=True, font=FONT_MONO)
    text(s, "EWMA z-scoring throws away level information that has predictive content for the panel-regression mean-reversion question. We disclose this openly. The BSI's payoff is in cross-firm threshold comparability for event detection (5/5 sensitivity), NOT in panel R². The panel and event-detection results test distinct claims.",
         M + Inches(0.5), Inches(5.8), SW - 2 * M - Inches(1.0), Inches(0.65),
         size=11, color=TEXT_DIM, bold=True, italic=False, line_spacing=1.45)
    footer(s, 20)


# ===============================================================================
# SLIDE 21 - Empirical results · SE-sensitivity table  (was slide 18)
# ===============================================================================
def s_results():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 20, "Empirical results · SE-sensitivity")
    slide_title(s, "Same coefficient.  Six standard-error (SE) specifications.  All reject H₀.",
                "Driscoll-Kraay survives at p = 0.007 — the strongest defensive result in the paper.")

    cols = ["SE estimator", "β̂", "SE", "t-stat", "p-value", "verdict"]
    rows = [
        ("Naive",                      "−0.082", "0.024", "−3.42", "0.001", "rejects H₀"),
        ("HC1 White",                  "−0.082", "0.029", "−2.83", "0.005", "rejects H₀"),
        ("heteroskedasticity- and autocorrelation-consistent (HAC) Newey-West (12)",        "−0.082", "0.031", "−2.65", "0.010", "rejects H₀"),
        ("Cluster (firm)",             "−0.082", "0.030", "−2.73", "0.007", "rejects H₀"),
        ("Two-way (firm × month)",     "−0.082", "0.034", "−2.41", "0.018", "rejects H₀"),
        ("Driscoll-Kraay  ★",          "−0.082", "0.030", "−2.71", "0.007", "rejects H₀"),
    ]
    y0 = Inches(2.95)
    rh = Inches(0.45)
    col_w = [Inches(3.5), Inches(1.2), Inches(1.2), Inches(1.4), Inches(1.4), Inches(2.6)]
    col_x = [M + Inches(0.4)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    for c_i, hdr in enumerate(cols):
        text(s, hdr, col_x[c_i], y0, col_w[c_i], rh,
             size=10, color=TEXT_FAINT, bold=True, font=FONT_MONO)
    hairline(s, M + Inches(0.4), y0 + Inches(0.42), SW - M - Inches(0.4), y0 + Inches(0.42))

    for r_i, row in enumerate(rows):
        y = y0 + (r_i + 1) * rh
        is_dk = "Driscoll" in row[0]
        for c_i, val in enumerate(row):
            if is_dk:
                col = ORANGE
            elif c_i == 5:
                col = GREEN
            elif c_i in (1, 2, 3):
                col = TEXT_DIM
            else:
                col = TEXT
            font = FONT_MONO
            size = 12 if c_i == 0 else 11
            text(s, val, col_x[c_i], y + Inches(0.05), col_w[c_i], rh,
                 size=size, color=col, font=font, bold=is_dk)

    text(s, "Driscoll-Kraay handles cross-sectional dependence and serial correlation simultaneously. Survival here ends the SE-robustness conversation.",
         M, Inches(6.4), SW - 2 * M, Inches(0.4),
         size=11, color=ORANGE_DEEP, italic=True, align=PP_ALIGN.CENTER)
    footer(s, 21)


# ===============================================================================
# SLIDE 14 - Future research
# ===============================================================================
def s_future():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 21, "Future research")
    slide_title(s, "Two open threads.",
                "Tier 2a fixed-income extension  ·  cross-asset BNPL ↔ credit-card contagion.")

    col_w = (SW - 2 * M - Inches(0.4)) / 2

    # LEFT — Tier 2a (the main thread)
    panel(s, M, Inches(2.85), col_w, Inches(3.0),
          fill=BG_PANEL, border=ORANGE, border_w=1.5)
    text(s, "THREAD 1  ·  TIER 2a  ·  fixed-income LHS",
         M + Inches(0.3), Inches(3.0), col_w, Inches(0.3),
         size=10, color=ORANGE, bold=True, font=FONT_MONO)
    text(s, "Δ Spreadₜ:ₜ₊₃₀  =  α + β · BSIₜ₋₁ + γ' Mₜ + εₜ",
         M + Inches(0.3), Inches(3.45), col_w - Inches(0.6), Inches(0.6),
         size=14, color=TEXT, bold=True, font=FONT_MONO, line_spacing=1.45)
    block(s, [
        {"text": "AFRMT junior I-spread, daily  ·  same six standard-error suite  ·  pre-registered",
         "size": 12, "color": TEXT_DIM, "bold": True, "italic": False, "space_before": 18, "line_spacing": 1.5},
        {"text": "Addresses the equity calendar-time α null directly.",
         "size": 12, "color": TEXT, "space_before": 8, "line_spacing": 1.5},
        {"text": "Data needed:  TRACE  ·  Markit CDS  ·  Intex trustee 10-D",
         "size": 11, "color": GOLD, "bold": True, "italic": False, "space_before": 12, "font": FONT_MONO},
    ], M + Inches(0.3), Inches(4.2), col_w - Inches(0.6), Inches(1.6))

    # RIGHT — cross-asset contagion (the second thread)
    panel(s, M + col_w + Inches(0.4), Inches(2.85), col_w, Inches(3.0),
          fill=BG_PANEL, border=PLUM, border_w=1.5)
    text(s, "THREAD 2  ·  CROSS-ASSET CONTAGION",
         M + col_w + Inches(0.7), Inches(3.0), col_w, Inches(0.3),
         size=10, color=PLUM, bold=True, font=FONT_MONO)
    text(s, "BNPL stress  ⇄  credit-card stress",
         M + col_w + Inches(0.7), Inches(3.45), col_w - Inches(0.6), Inches(0.5),
         size=16, color=TEXT, bold=True, font=FONT_HEAD)
    block(s, [
        {"text": "Does BNPL distress lead credit-card delinquencies — or vice versa?",
         "size": 12, "color": TEXT, "bold": True, "space_before": 14, "line_spacing": 1.5},
        {"text": "BNPL borrowers are 2.4× more likely to be subprime  →  same wallet, two products, one stress signal?",
         "size": 11, "color": TEXT_DIM, "italic": False, "space_before": 8, "line_spacing": 1.5},
        {"text": "Methodology: same panel regression, LHS = card-monoline charge-offs; RHS = BNPL-firm BSI.",
         "size": 11, "color": PLUM, "bold": True, "italic": False, "space_before": 12, "font": FONT_MONO, "line_spacing": 1.4},
    ], M + col_w + Inches(0.7), Inches(4.0), col_w - Inches(0.6), Inches(1.8))

    # Bottom note
    text(s, "Both threads use the same engine. Tier 2a is data-bound; cross-asset is methodology-design-bound.",
         M, Inches(6.2), SW - 2 * M, Inches(0.4),
         size=12, color=TEXT_DIM, bold=True, italic=False, align=PP_ALIGN.CENTER)
    footer(s, 22)


# ===============================================================================
# SLIDE 15 - Contribution
# ===============================================================================
def s_contribution():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 22, "Contribution")
    slide_title(s, "Three things this paper contributes.",
                "What is established, what is new, what is deliberately not claimed.")

    col_w = (SW - 2 * M - Inches(0.4)) / 3
    y0 = Inches(2.95)

    # 1. ESTABLISHED
    panel(s, M, y0, col_w, Inches(3.5),
          fill=BG_PANEL, border=GREEN, border_w=1.2)
    text(s, "ESTABLISHED",
         M + Inches(0.3), y0 + Inches(0.2), col_w, Inches(0.3),
         size=10, color=GREEN, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "Generalises to non-bank consumer credit",
         "size": 13, "color": TEXT, "bold": True, "space_before": 10, "line_spacing": 1.4},
        {"text": "BNPL · subprime auto · marketplace lending.",
         "size": 11, "color": TEXT_DIM, "space_before": 4, "line_spacing": 1.5},

        {"text": "Equity LHS panel result is robust",
         "size": 13, "color": TEXT, "bold": True, "space_before": 14, "line_spacing": 1.4},
        {"text": "Driscoll-Kraay  p = 0.007.",
         "size": 11, "color": TEXT_DIM, "space_before": 4},
    ], M + Inches(0.3), y0 + Inches(0.6), col_w - Inches(0.6), Inches(2.8))

    # 2. NEW
    panel(s, M + col_w + Inches(0.2), y0, col_w, Inches(3.5),
          fill=BG_PANEL, border=ORANGE, border_w=1.2)
    text(s, "METHODOLOGICALLY NEW",
         M + col_w + Inches(0.5), y0 + Inches(0.2), col_w, Inches(0.3),
         size=10, color=ORANGE, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "Multi-pillar pre-registered composite",
         "size": 13, "color": TEXT, "bold": True, "space_before": 10, "line_spacing": 1.4},
        {"text": "First fusion of CFPB with seven additional pillars under hash-locked weights.",
         "size": 11, "color": TEXT_DIM, "space_before": 4, "line_spacing": 1.5},

        {"text": "Six-estimator robustness suite",
         "size": 13, "color": TEXT, "bold": True, "space_before": 14, "line_spacing": 1.4},
        {"text": "More demanding than typical alt-data work.",
         "size": 11, "color": TEXT_DIM, "space_before": 4, "line_spacing": 1.5},
    ], M + col_w + Inches(0.5), y0 + Inches(0.6), col_w - Inches(0.6), Inches(2.8))

    # 3. NOT CLAIMED
    panel(s, M + 2 * (col_w + Inches(0.2)) - Inches(0.2), y0, col_w, Inches(3.5),
          fill=BG_PANEL, border=GOLD, border_w=1.2)
    text(s, "NOT YET CLAIMED",
         M + 2 * (col_w + Inches(0.2)) + Inches(0.1), y0 + Inches(0.2), col_w, Inches(0.3),
         size=10, color=GOLD, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "Equity calendar-time alpha",
         "size": 13, "color": TEXT, "bold": True, "space_before": 10, "line_spacing": 1.4},
        {"text": "Honest null  t = 0.08.",
         "size": 11, "color": TEXT_DIM, "space_before": 4},

        {"text": "Fixed-income alpha",
         "size": 13, "color": TEXT, "bold": True, "space_before": 14, "line_spacing": 1.4},
        {"text": "Tier 2a positioned, not yet executed.",
         "size": 11, "color": TEXT_DIM, "space_before": 4, "line_spacing": 1.5},
    ], M + 2 * (col_w + Inches(0.2)) + Inches(0.1), y0 + Inches(0.6), col_w - Inches(0.6), Inches(2.8))

    text(s, "The discipline of not overclaiming is what makes the established claims defensible.",
         M, Inches(6.7), SW - 2 * M, Inches(0.3),
         size=12, color=TEXT, italic=True, align=PP_ALIGN.CENTER)
    footer(s, 23)


# ===============================================================================
# SLIDE 16 - False positive + total return swap (TRS) opportunity bridge
# ===============================================================================
def s_fp_trs():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 23, "Bridge to the pod")
    slide_title(s, "One known false positive.  One open opportunity.",
                "Sezzle shows where the methodology mis-fires; junior-tranche TRS shows where it pays off.")

    col_w = (SW - 2 * M - Inches(0.4)) / 2

    # FP — SEZL
    panel(s, M, Inches(2.85), col_w, Inches(3.6),
          fill=BG_PANEL, border=RED_TERRA, border_w=1.2)
    text(s, "FALSE POSITIVE  ·  SEZL 2023-2025",
         M + Inches(0.3), Inches(3.0), col_w, Inches(0.3),
         size=10, color=RED_TERRA, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "BSI fired repeatedly. Equity went +600%.",
         "size": 13, "color": TEXT, "bold": True, "space_before": 12, "line_spacing": 1.4},
        {"text": "Why  ·  growth-stage issuer; complaint volume scaled with active-customer expansion.",
         "size": 11, "color": TEXT_DIM, "space_before": 10, "line_spacing": 1.5},
        {"text": "Numerator rose mechanically. Per-customer rate stayed stable.",
         "size": 11, "color": TEXT_DIM, "space_before": 6, "line_spacing": 1.5},
        {"text": "Documented in §13.3. Denominator-normalisation refinement addresses it.",
         "size": 11, "color": GOLD, "italic": True, "space_before": 12, "font": FONT_MONO},
    ], M + Inches(0.3), Inches(3.4), col_w - Inches(0.6), Inches(3.0))

    # OPP — TRS
    panel(s, M + col_w + Inches(0.4), Inches(2.85), col_w, Inches(3.6),
          fill=BG_PANEL, border=GREEN, border_w=1.2)
    text(s, "OPPORTUNITY  ·  JUNIOR TRANCHE TRS",
         M + col_w + Inches(0.7), Inches(3.0), col_w, Inches(0.3),
         size=10, color=GREEN, bold=True, font=FONT_MONO)
    block(s, [
        {"text": "The instrument structure dodges the failure modes.",
         "size": 13, "color": TEXT, "bold": True, "space_before": 12, "line_spacing": 1.4},
        {"text": "Tranche cashflow doesn't care about customer growth — filters the SEZL case.",
         "size": 11, "color": TEXT_DIM, "space_before": 10, "line_spacing": 1.5},
        {"text": "Spreads price default risk directly — addresses the equity-α null.",
         "size": 11, "color": TEXT_DIM, "space_before": 6, "line_spacing": 1.5},
        {"text": "Methodology and trade vehicle are co-designed.",
         "size": 11, "color": GREEN, "italic": True, "space_before": 12, "font": FONT_MONO},
    ], M + col_w + Inches(0.7), Inches(3.4), col_w - Inches(0.6), Inches(3.0))

    footer(s, 24)


# ===============================================================================
# SLIDE 17 - POD · Case Study (full-bleed screenshot)
# ===============================================================================
def s_pod_case():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 24, "Pod · case study (CVNA)")

    img_x = M
    img_y = Inches(1.2)
    img_w = SW - 2 * M
    img_h = Inches(5.5)

    if CASE_STUDY_PNG.exists():
        s.shapes.add_picture(str(CASE_STUDY_PNG), img_x, img_y, width=img_w, height=img_h)
    else:
        panel(s, img_x, img_y, img_w, img_h, fill=BG_PANEL)
        text(s, "[ /case-study screenshot — capture instructions in build log ]",
             img_x, img_y + Inches(2.4), img_w, Inches(0.6),
             size=14, color=TEXT_DIM, italic=True,
             font=FONT_MONO, align=PP_ALIGN.CENTER)
        text(s, "Save to  screenshots/case_study.png  and rebuild.",
             img_x, img_y + Inches(3.0), img_w, Inches(0.4),
             size=11, color=TEXT_FAINT, italic=True, align=PP_ALIGN.CENTER)

    text(s, "/case-study  ·  CVNA 2022 deep-dive  ·  +96.3% short, 540d hold, 5/5 gates",
         M, Inches(6.85), SW - 2 * M, Inches(0.3),
         size=12, color=TEXT, bold=True, italic=False, align=PP_ALIGN.CENTER)
    footer(s, 25)


# ===============================================================================
# SLIDE 18 - POD · CVNA math + playbook (two screenshots)
# ===============================================================================
def s_pod_math_playbook():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 25, "Pod · case study (CVNA · trade math + playbook)")

    img_y = Inches(1.5)
    gap   = Inches(0.25)
    half_w = (SW - 2 * M - gap) / 2  # column width — height auto-computes from aspect ratio

    # LEFT — math (only set WIDTH; height preserves aspect ratio of original PNG)
    if CVNA_MATH_PNG.exists():
        from PIL import Image as _PILImg
        with _PILImg.open(CVNA_MATH_PNG) as im:
            iw, ih = im.size
        # Compute height that preserves aspect at half-column width
        nat_h_emu = int(half_w * ih / iw)
        # Centre vertically in available space
        avail_h = Inches(5.4)
        y_off = img_y + (avail_h - nat_h_emu) / 2 if nat_h_emu < avail_h else img_y
        s.shapes.add_picture(str(CVNA_MATH_PNG), M, y_off, width=half_w)
    else:
        panel(s, M, img_y, half_w, Inches(5.0), fill=BG_PANEL)
        text(s, "[ cvna_math.png ]",
             M, img_y + Inches(2.2), half_w, Inches(0.5),
             size=14, color=TEXT_DIM, bold=True, italic=True,
             font=FONT_MONO, align=PP_ALIGN.CENTER)

    # RIGHT — playbook (same aspect-preserving treatment)
    right_x = M + half_w + gap
    if CVNA_PLAYBOOK_PNG.exists():
        from PIL import Image as _PILImg
        with _PILImg.open(CVNA_PLAYBOOK_PNG) as im:
            iw, ih = im.size
        nat_h_emu = int(half_w * ih / iw)
        avail_h = Inches(5.4)
        y_off = img_y + (avail_h - nat_h_emu) / 2 if nat_h_emu < avail_h else img_y
        s.shapes.add_picture(str(CVNA_PLAYBOOK_PNG), right_x, y_off, width=half_w)
    else:
        panel(s, right_x, img_y, half_w, Inches(5.0), fill=BG_PANEL)
        text(s, "[ cvna_playbook.png ]",
             right_x, img_y + Inches(2.2), half_w, Inches(0.5),
             size=14, color=TEXT_DIM, bold=True, italic=True,
             font=FONT_MONO, align=PP_ALIGN.CENTER)

    text(s, "Trade math (left)  ·  Three-leg playbook (right)  ·  short equity + short credit + long-pod recovery",
         M, Inches(7.0), SW - 2 * M, Inches(0.3),
         size=11, color=TEXT_DIM, bold=True, italic=False, align=PP_ALIGN.CENTER)
    footer(s, 26)


# ===============================================================================
# SLIDE 19 - POD · Live Pod (full-bleed screenshot)
# ===============================================================================
def s_pod_live():
    s = prs.slides.add_slide(BLANK); bg(s)
    kicker(s, 26, "Pod · live")

    img_x = M
    img_y = Inches(1.2)
    img_w = SW - 2 * M
    img_h = Inches(5.5)

    if LIVE_POD_PNG.exists():
        s.shapes.add_picture(str(LIVE_POD_PNG), img_x, img_y, width=img_w, height=img_h)
    else:
        panel(s, img_x, img_y, img_w, img_h, fill=BG_PANEL)
        text(s, "[ /live screenshot — capture instructions in build log ]",
             img_x, img_y + Inches(2.4), img_w, Inches(0.6),
             size=14, color=TEXT_DIM, italic=True,
             font=FONT_MONO, align=PP_ALIGN.CENTER)
        text(s, "Save to  screenshots/live_pod.png  and rebuild.",
             img_x, img_y + Inches(3.0), img_w, Inches(0.4),
             size=11, color=TEXT_FAINT, italic=True, align=PP_ALIGN.CENTER)

    text(s, "/live  ·  4-stage pipeline  ·  5-gate evaluation  ·  technical-override layer",
         M, Inches(6.85), SW - 2 * M, Inches(0.3),
         size=12, color=TEXT, bold=True, italic=False, align=PP_ALIGN.CENTER)
    footer(s, 27)


# ===============================================================================
# SLIDE 20 - Q&A
# ===============================================================================
def s_qa():
    s = prs.slides.add_slide(BLANK); bg(s)

    # Logo top-centre, smaller to leave room for thank-you + acknowledgements + Q&A
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), SW / 2 - Inches(0.42), Inches(0.55),
                             height=Inches(0.7))

    orange_bar(s, M, Inches(1.55), SW - 2 * M, Inches(0.04))

    # THANK YOU heading (calligraphic Cambria, large)
    text(s, "Thank you",
         M, Inches(1.85), SW - 2 * M, Inches(1.0),
         size=56, color=TEXT, bold=True, font=FONT_DISP, align=PP_ALIGN.CENTER)

    # Acknowledgements panel
    ack_x = M + Inches(2.0)
    ack_w = SW - 2 * M - Inches(4.0)
    ack_y = Inches(3.05)
    panel(s, ack_x, ack_y, ack_w, Inches(2.1),
          fill=BG_PANEL, border=ORANGE, border_w=1.5)
    text(s, "ACKNOWLEDGEMENTS",
         ack_x, ack_y + Inches(0.18), ack_w, Inches(0.3),
         size=11, color=ORANGE, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # Four names — each on its own line for visual breath
    names = [
        "Claude",
        "Professor Tony Zhang",
        "Professor Martin Widdicks",
        "Professor Neil Pearson",
    ]
    name_y0 = ack_y + Inches(0.6)
    name_h  = Inches(0.32)
    for i, nm in enumerate(names):
        text(s, nm, ack_x, name_y0 + i * name_h, ack_w, name_h,
             size=15, color=TEXT, bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER)

    # Q&A line
    text(s, "Q & A",
         M, Inches(5.5), SW - 2 * M, Inches(0.6),
         size=32, color=ORANGE_DEEP, bold=True, font=FONT_DISP, align=PP_ALIGN.CENTER)
    text(s, "questions, comments, criticisms welcome",
         M, Inches(6.15), SW - 2 * M, Inches(0.4),
         size=13, color=TEXT_DIM, bold=True, italic=False, align=PP_ALIGN.CENTER)

    # Bottom meta
    text(s, "Siddharth Verma  ·  FIN 580  ·  University of Illinois Urbana-Champaign  ·  Spring 2026",
         M, SH - Inches(0.45), SW - 2 * M, Inches(0.3),
         size=11, color=TEXT_FAINT, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)


# ===============================================================================
# Build
# ===============================================================================
if __name__ == "__main__":
    s_title()
    s_motivation()
    s_what_is_bsi()        # slide 3 — visual definition of BSI
    s_neural_diagram()     # NEW slide 4 — 2D neural-style architecture diagram
    s_hypothesis()
    s_cascade()
    s_data()
    s_data_treatment()
    s_regression()
    s_analysis()
    s_findings()
    s_table_events()       # NEW slide 10 — canonical events table (paper Table tab:sens)
    s_event_calendar()     # NEW slide 12 — event calendar timeline
    s_case_cvna()
    s_case_minis()
    s_caveats()
    s_table_panel()        # slide 17 — panel regression coefficients
    s_table_normalised()   # NEW slide 18 — normalised BSI vs revenue YoY
    s_table_granger()      # NEW slide 19 — Granger F-tests by lag
    s_table_baseline()     # NEW slide 20 — baseline comparison (honest disclosure)
    s_results()
    s_future()
    s_contribution()
    s_fp_trs()
    s_pod_case()
    s_pod_math_playbook()
    s_pod_live()
    s_qa()

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"OK · wrote {OUT}")
    print(f"   slides: {len(prs.slides)}")
    import os
    print(f"   bytes : {os.path.getsize(OUT):,}")
    print()
    if not (CASE_STUDY_PNG.exists() and LIVE_POD_PNG.exists()):
        missing = []
        if not CASE_STUDY_PNG.exists(): missing.append(CASE_STUDY_PNG.name)
        if not LIVE_POD_PNG.exists():   missing.append(LIVE_POD_PNG.name)
        print(f"   placeholders for: {', '.join(missing)}")
        print(f"   capture and save to {SCREENSHOT_DIR}, then rerun")
    else:
        print("   both screenshots embedded")
